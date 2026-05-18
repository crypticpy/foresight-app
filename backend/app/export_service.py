"""
Export Service for Intelligence Cards and Workstream Reports.

This service handles generation of export files in multiple formats:
- PDF: Using ReportLab for professional document generation
- PowerPoint: Using python-pptx for presentation slides
- CSV: Using pandas for tabular data export

Chart Generation:
- Score radar charts showing all dimensions
- Score bar charts for individual scores
- Pillar distribution charts for workstream reports
- All charts use matplotlib with 'Agg' backend (non-GUI)

Usage:
    export_service = ExportService(supabase_client)
    pdf_path = await export_service.generate_pdf(card_data)
    pptx_path = await export_service.generate_pptx(card_data)
    csv_content = await export_service.generate_csv(card_data)
"""

import logging
import re
import tempfile
from datetime import datetime, timezone
from typing import Dict, List, Optional, Any, Tuple

import matplotlib

matplotlib.use("Agg")  # Non-GUI backend - must be set before importing pyplot
import matplotlib.pyplot as plt
import numpy as np

from app.openai_provider import get_chat_deployment

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ReportLab imports for PDF generation
from reportlab.lib import colors as rl_colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    Image as RLImage,
    PageBreak,
    HRFlowable,
)

from supabase import Client

from .models.export import (
    CardExportData,
    ExportFormat,
)

# Import classification definitions from gamma_service for backup slides
from .gamma_service import (
    PILLAR_DEFINITIONS,
    HORIZON_DEFINITIONS,
    STAGE_DEFINITIONS,
)

# Re-exported helpers (definitions live in app.export package). Imported here so
# legacy callers can continue to reach them via app.export_service.
from .export.branding import (
    COA_BRAND_COLORS,
    FORESIGHT_COLORS,
    PDF_COLORS,
)
from .export import briefs as _briefs
from .export import cards as _cards
from .export import charts as _charts
from .export import csv_export as _csv_export
from .export import data_access as _data_access
from .export import utils as _utils
from .export import pptx as _pptx_components
from .export import workstreams as _workstreams
from .export.charts import CHART_DPI
from .export.pptx import (
    PPTX_MARGIN,
    PPTX_SLIDE_HEIGHT,
    PPTX_SLIDE_WIDTH,
)

logger = logging.getLogger(__name__)


# ============================================================================
# Dimensional constants used by ExportService methods
# ============================================================================

# PDF settings
PDF_PAGE_SIZE = letter
PDF_MARGIN = 0.75 * inch
PDF_TITLE_FONT_SIZE = 24
PDF_HEADING_FONT_SIZE = 14
PDF_BODY_FONT_SIZE = 11
PDF_SMALL_FONT_SIZE = 9
PDF_CHART_WIDTH = 5.5 * inch
PDF_CHART_HEIGHT = 4 * inch


# ============================================================================
# Export Service
# ============================================================================


class ExportService:
    """
    Service for generating export files from intelligence cards and workstreams.

    Supports PDF, PowerPoint, and CSV export formats with embedded visualizations.
    Follows the service class pattern from research_service.py.
    """

    def __init__(self, supabase: Client):
        """
        Initialize the ExportService.

        Args:
            supabase: Supabase client for database queries
        """
        self.supabase = supabase
        logger.info("ExportService initialized")

    # ========================================================================
    # Chart Generation Methods (facade — implementations live in export.charts)
    # ========================================================================

    def generate_score_chart(
        self, card_data: CardExportData, chart_type: str = "bar", dpi: int = CHART_DPI
    ) -> Optional[str]:
        return _charts.generate_score_chart(card_data, chart_type, dpi)

    def _generate_bar_chart(
        self, scores: Dict[str, int], title: str, dpi: int
    ) -> str:
        return _charts.generate_bar_chart(scores, title, dpi)

    def _generate_radar_chart(
        self, scores: Dict[str, int], title: str, dpi: int
    ) -> str:
        return _charts.generate_radar_chart(scores, title, dpi)

    def generate_pillar_distribution_chart(
        self,
        pillar_counts: Dict[str, int],
        title: str = "Pillar Distribution",
        dpi: int = CHART_DPI,
    ) -> Optional[str]:
        return _charts.generate_pillar_distribution_chart(pillar_counts, title, dpi)

    def generate_horizon_distribution_chart(
        self,
        horizon_counts: Dict[str, int],
        title: str = "Horizon Distribution",
        dpi: int = CHART_DPI,
    ) -> Optional[str]:
        return _charts.generate_horizon_distribution_chart(horizon_counts, title, dpi)

    # ========================================================================
    # PDF Generation Methods
    # ========================================================================

    def _create_pdf_header(
        self, card_data: CardExportData, styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF header elements for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the header
        """
        elements = []

        # Title
        title = Paragraph(card_data.name, styles["Title"])
        elements.extend((title, Spacer(1, 6)))
        # Horizontal rule
        elements.append(
            HRFlowable(
                width="100%",
                thickness=2,
                color=PDF_COLORS["primary"],
                spaceBefore=6,
                spaceAfter=12,
            )
        )

        # Metadata badges (pillar, horizon, stage)
        badge_data = []
        if card_data.pillar_name or card_data.pillar_id:
            badge_data.append(("Pillar", card_data.pillar_name or card_data.pillar_id))
        if card_data.horizon:
            badge_data.append(("Horizon", card_data.horizon))
        if card_data.stage_name or card_data.stage_id:
            badge_data.append(("Stage", card_data.stage_name or card_data.stage_id))
        if card_data.goal_name or card_data.goal_id:
            badge_data.append(("Goal", card_data.goal_name or card_data.goal_id))

        if badge_data:
            badge_table_data = [[]]
            for label, value in badge_data:
                badge_text = f"<b>{label}:</b> {value}"
                badge_table_data[0].append(Paragraph(badge_text, styles["Small"]))

            badge_table = Table(
                badge_table_data, colWidths=[1.5 * inch] * len(badge_data)
            )
            badge_table.setStyle(
                TableStyle(
                    [
                        ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("TOPPADDING", (0, 0), (-1, -1), 4),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                    ]
                )
            )
            elements.append(badge_table)
            elements.append(Spacer(1, 12))

        return elements

    def _create_pdf_summary_section(
        self, card_data: CardExportData, styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF summary section for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the summary
        """
        elements = [Paragraph("Classification Overview", styles["Heading1"])]

        classification_items = []
        if card_data.pillar_name or card_data.pillar_id:
            classification_items.append(
                f"<b>Strategic Pillar:</b> {card_data.pillar_name or card_data.pillar_id}"
            )
        if card_data.goal_name or card_data.goal_id:
            classification_items.append(
                f"<b>Strategic Goal:</b> {card_data.goal_name or card_data.goal_id}"
            )
        if card_data.horizon:
            horizon_desc = {
                "H1": "Near-term (0-2 years)",
                "H2": "Mid-term (2-5 years)",
                "H3": "Long-term (5+ years)",
            }
            classification_items.append(
                f"<b>Time Horizon:</b> {card_data.horizon} - {horizon_desc.get(card_data.horizon, '')}"
            )
        if card_data.stage_name or card_data.stage_id:
            classification_items.append(
                f"<b>Maturity Stage:</b> {card_data.stage_name or card_data.stage_id}"
            )

        if classification_items:
            for item in classification_items:
                elements.append(Paragraph(item, styles["Body"]))
            elements.append(Spacer(1, 12))

        # Executive Summary
        if card_data.summary:
            elements.append(Paragraph("Executive Summary", styles["Heading1"]))
            elements.append(Paragraph(card_data.summary, styles["Body"]))
            elements.append(Spacer(1, 12))

        # Full Description
        if card_data.description:
            elements.append(Paragraph("Detailed Analysis", styles["Heading1"]))
            # Truncate very long descriptions
            description = card_data.description
            if len(description) > 3000:
                description = description[:3000] + "... [truncated]"
            elements.append(Paragraph(description, styles["Body"]))
            elements.append(Spacer(1, 12))

        return elements

    def _create_pdf_research_section(
        self, card_data: CardExportData, styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF section for deep research report.

        Args:
            card_data: Card data containing research report
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the research section
        """
        elements = []

        if not card_data.deep_research_report:
            return elements

        elements.extend(
            (
                PageBreak(),
                Paragraph("Strategic Intelligence Report", styles["Title"]),
                Spacer(1, 6),
                Paragraph(
                    "<i>The following strategic intelligence report was generated through deep research analysis, "
                    "synthesizing multiple sources to provide actionable insights for decision-makers.</i>",
                    styles["Small"],
                ),
                Spacer(1, 12),
                HRFlowable(
                    width="100%",
                    thickness=1,
                    color=PDF_COLORS["secondary"],
                    spaceBefore=6,
                    spaceAfter=12,
                ),
            )
        )
        # Parse and render the markdown report
        # Simple markdown-to-PDF conversion for common elements
        report = card_data.deep_research_report

        # Truncate if extremely long
        if len(report) > 15000:
            report = (
                report[:15000]
                + "\n\n... [Report truncated for PDF export. View full report in the application.]"
            )

        # Process the report line by line
        lines = report.split("\n")
        for line in lines:
            line = line.strip()
            if not line:
                elements.append(Spacer(1, 6))
            elif line.startswith("# "):
                # Main heading
                elements.append(Paragraph(line[2:], styles["Heading1"]))
            elif line.startswith("## "):
                # Section heading
                elements.append(Paragraph(line[3:], styles["Heading2"]))
            elif line.startswith("### "):
                # Subsection heading
                text = f"<b>{line[4:]}</b>"
                elements.append(Paragraph(text, styles["Body"]))
            elif line.startswith("- ") or line.startswith("* "):
                # Bullet point
                text = f"• {line[2:]}"
                elements.append(Paragraph(text, styles["Body"]))
            elif line.startswith("**") and line.endswith("**"):
                # Bold text
                text = f"<b>{line[2:-2]}</b>"
                elements.append(Paragraph(text, styles["Body"]))
            else:
                # Regular paragraph
                # Escape any XML-unsafe characters
                safe_line = (
                    line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                )
                elements.append(Paragraph(safe_line, styles["Body"]))

        return elements

    def _create_pdf_scores_table(
        self, card_data: CardExportData, styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF scores table for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the scores section
        """
        elements = [Paragraph("Scores", styles["Heading1"])]

        # Build scores table
        scores = card_data.get_all_scores()
        table_data = [["Metric", "Score", "Rating"]]

        for name, score in scores.items():
            score_display = self.format_score_display(score)

            # Determine rating based on score
            if score is None:
                rating = "Not Scored"
            elif score >= 80:
                rating = "Excellent"
            elif score >= 60:
                rating = "Good"
            elif score >= 40:
                rating = "Fair"
            else:
                rating = "Low"

            table_data.append([name, score_display, rating])

        # Create table with styling
        table = Table(table_data, colWidths=[2 * inch, 1 * inch, 1.5 * inch])
        table.setStyle(
            TableStyle(
                [
                    # Header styling
                    ("BACKGROUND", (0, 0), (-1, 0), PDF_COLORS["primary"]),
                    ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, 0), PDF_BODY_FONT_SIZE),
                    ("ALIGN", (0, 0), (-1, 0), "CENTER"),
                    # Body styling
                    ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                    ("FONTSIZE", (0, 1), (-1, -1), PDF_BODY_FONT_SIZE),
                    ("ALIGN", (1, 1), (1, -1), "CENTER"),  # Center score column
                    ("ALIGN", (2, 1), (2, -1), "CENTER"),  # Center rating column
                    # Alternating row colors
                    (
                        "ROWBACKGROUNDS",
                        (0, 1),
                        (-1, -1),
                        [rl_colors.white, PDF_COLORS["light"]],
                    ),
                    # Grid
                    ("GRID", (0, 0), (-1, -1), 0.5, PDF_COLORS["light"]),
                    ("BOX", (0, 0), (-1, -1), 1, PDF_COLORS["primary"]),
                    # Padding
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                    ("LEFTPADDING", (0, 0), (-1, -1), 8),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                ]
            )
        )

        elements.append(table)
        elements.append(Spacer(1, 18))

        return elements

    def _create_pdf_chart_section(
        self, card_data: CardExportData, styles: Dict[str, ParagraphStyle]
    ) -> Tuple[List[Any], List[str]]:
        """
        Create PDF chart section for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            Tuple of (list of flowable elements, list of temp file paths to clean up)
        """
        elements = []
        temp_files = []

        if chart_path := self.generate_score_chart(card_data, chart_type="bar"):
            temp_files.append(chart_path)
            elements.append(Paragraph("Score Visualization", styles["Heading1"]))

            try:
                img = RLImage(
                    chart_path, width=PDF_CHART_WIDTH, height=PDF_CHART_HEIGHT
                )
                elements.extend((img, Spacer(1, 12)))
            except Exception as e:
                logger.warning(f"Failed to add chart image to PDF: {e}")

        return elements, temp_files

    def _create_pdf_footer(
        self, card_data: CardExportData, styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF footer elements for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the footer
        """
        elements = [Spacer(1, 24)]

        elements.append(
            HRFlowable(
                width="100%",
                thickness=1,
                color=PDF_COLORS["light"],
                spaceBefore=6,
                spaceAfter=6,
            )
        )

        # Metadata footer
        footer_parts = []
        if card_data.created_at:
            footer_parts.append(f"Created: {card_data.created_at.strftime('%Y-%m-%d')}")
        if card_data.updated_at:
            footer_parts.append(f"Updated: {card_data.updated_at.strftime('%Y-%m-%d')}")
        footer_parts.append(
            f"Export Date: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}"
        )

        footer_text = " | ".join(footer_parts)
        elements.append(Paragraph(footer_text, styles["Small"]))

        # Branding
        elements.append(Spacer(1, 6))
        elements.append(
            Paragraph("Generated by Foresight Intelligence Platform", styles["Small"])
        )

        return elements

    async def generate_pdf(
        self, card_data: CardExportData, include_charts: bool = True
    ) -> str:
        return await _cards.generate_card_pdf(card_data, include_charts)

    async def generate_workstream_pdf(
        self, workstream_id: str, include_charts: bool = True, max_cards: int = 50
    ) -> str:
        return await _workstreams.generate_workstream_pdf(
            self.supabase, workstream_id, include_charts, max_cards
        )

    # ========================================================================
    # Utility Methods (facade — implementations live in export.utils + export.data_access)
    # ========================================================================

    def cleanup_temp_files(self, file_paths: List[str]) -> None:
        _utils.cleanup_temp_files(file_paths)

    async def get_card_data(self, card_id: str) -> Optional[CardExportData]:
        return await _data_access.get_card_data(self.supabase, card_id)

    async def get_workstream_cards(
        self, workstream_id: str, max_cards: int = 50
    ) -> Tuple[Optional[Dict[str, Any]], List[CardExportData]]:
        return await _data_access.get_workstream_cards(
            self.supabase, workstream_id, max_cards
        )

    def format_score_display(self, score: Optional[int]) -> str:
        return _utils.format_score_display(score)

    def get_content_type(self, format: ExportFormat) -> str:
        return _utils.get_content_type(format)

    def generate_filename(self, name: str, format: ExportFormat) -> str:
        return _utils.generate_filename(name, format)

    # ========================================================================
    # CSV Export Methods
    # ========================================================================

    async def generate_csv(
        self,
        card_data: CardExportData,
    ) -> str:
        return await _csv_export.generate_csv(card_data)

    async def generate_csv_multi(
        self,
        cards: List[CardExportData],
    ) -> str:
        return await _csv_export.generate_csv_multi(cards)

    def _generate_empty_csv(self) -> str:
        return _csv_export.generate_empty_csv()

    # ========================================================================
    # PowerPoint Export Methods
    # ========================================================================

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        return _pptx_components.hex_to_rgb(hex_color)

    def _add_pptx_header(self, slide, include_logo: bool = True) -> None:
        _pptx_components.add_pptx_header(slide, include_logo)

    def _add_pptx_footer(self, slide, include_ai_disclosure: bool = True) -> None:
        _pptx_components.add_pptx_footer(slide, include_ai_disclosure)

    def _add_title_slide(
        self, prs: Presentation, title: str, subtitle: Optional[str] = None
    ) -> None:
        _pptx_components.add_title_slide(prs, title, subtitle)

    def _add_content_slide(
        self,
        prs: Presentation,
        title: str,
        content_items: List[Tuple[str, str]],
        chart_path: Optional[str] = None,
    ) -> None:
        _pptx_components.add_content_slide(prs, title, content_items, chart_path)

    def _add_scores_slide(
        self,
        prs: Presentation,
        card_data: CardExportData,
        chart_path: Optional[str] = None,
    ) -> None:
        _pptx_components.add_scores_slide(prs, card_data, chart_path)

    def _add_description_slide(
        self, prs: Presentation, title: str, description: Optional[str]
    ) -> None:
        _pptx_components.add_description_slide(prs, title, description)

    async def generate_pptx(
        self,
        card_data: CardExportData,
        include_charts: bool = True,
        include_description: bool = True,
    ) -> str:
        return await _cards.generate_card_pptx(
            card_data, include_charts, include_description
        )

    async def generate_workstream_pptx(
        self,
        workstream: Dict[str, Any],
        cards: List[CardExportData],
        include_charts: bool = True,
        include_card_details: bool = True,
    ) -> str:
        return await _workstreams.generate_workstream_pptx(
            workstream, cards, include_charts, include_card_details
        )

    # ========================================================================
    # Executive Brief Export Methods
    # ========================================================================

    async def generate_brief_pdf(
        self,
        brief_title: str,
        card_name: str,
        executive_summary: str,
        content_markdown: str,
        generated_at: Optional[datetime] = None,
        version: Optional[int] = None,
    ) -> str:
        """Generate a basic PDF export for an executive brief."""
        return await _briefs.generate_brief_pdf(
            brief_title=brief_title,
            card_name=card_name,
            executive_summary=executive_summary,
            content_markdown=content_markdown,
            generated_at=generated_at,
            version=version,
        )

    async def generate_professional_brief_pdf(
        self,
        brief_title: str,
        card_name: str,
        executive_summary: str,
        content_markdown: str,
        generated_at: Optional[datetime] = None,
        version: Optional[int] = None,
        classification: Optional[Dict[str, str]] = None,
    ) -> str:
        """Generate a branded, mayor-ready PDF for an executive brief."""
        return await _briefs.generate_professional_brief_pdf(
            brief_title=brief_title,
            card_name=card_name,
            executive_summary=executive_summary,
            content_markdown=content_markdown,
            generated_at=generated_at,
            version=version,
            classification=classification,
        )

    async def generate_chat_response_pdf(
        self,
        title: str,
        question: str,
        response_content: str,
        citations: Optional[List[Dict[str, Any]]] = None,
        metadata: Optional[Dict[str, Any]] = None,
        scope: Optional[str] = None,
        scope_context: Optional[str] = None,
    ) -> str:
        """Generate a professional PDF export for a chat response."""
        return await _briefs.generate_chat_response_pdf(
            title=title,
            question=question,
            response_content=response_content,
            citations=citations,
            metadata=metadata,
            scope=scope,
            scope_context=scope_context,
        )

    async def generate_brief_pptx(
        self,
        brief_title: str,
        card_name: str,
        executive_summary: str,
        content_markdown: str,
        generated_at: Optional[datetime] = None,
        version: Optional[int] = None,
        classification: Optional[Dict[str, str]] = None,
        use_gamma: bool = True,
    ) -> str:
        """Generate a PowerPoint presentation for an executive brief."""
        return await _briefs.generate_brief_pptx(
            brief_title=brief_title,
            card_name=card_name,
            executive_summary=executive_summary,
            content_markdown=content_markdown,
            generated_at=generated_at,
            version=version,
            classification=classification,
            use_gamma=use_gamma,
        )

    # =========================================================================
    # Portfolio Export (Bulk Brief Export)
    # =========================================================================
    # The two helpers below are still used by portfolio export methods that
    # remain on this class. They delegate to app.export.briefs.pptx so the
    # implementation stays in one place. They will be inlined into the
    # portfolios module when portfolio export is extracted.

    def _add_smart_content_slide(
        self, prs: Presentation, title: str, content: str, max_chars: int = 1000
    ) -> None:
        _briefs.add_smart_content_slide(prs, title, content, max_chars)

    def _add_ai_disclosure_slide(self, prs: Presentation) -> None:
        _briefs.add_ai_disclosure_slide(prs)

    def _extract_key_takeaways(self, brief_markdown: str) -> List[str]:
        """
        Extract key takeaways from brief markdown content.

        Looks for sections like "Key Takeaways", "Key Findings", "Key Implications",
        "What This Means", bullet points, or numbered lists.
        """

        takeaways = []

        # Try to find key sections
        key_section_patterns = [
            r"(?:##?\s*)?(?:Key\s+)?(?:Takeaways?|Findings?|Implications?|Insights?)[\s:]*\n((?:[-•*]\s*.+\n?)+)",
            r"(?:##?\s*)?What\s+This\s+Means[^:]*:?\s*\n((?:[-•*]\s*.+\n?)+)",
            r"(?:##?\s*)?Strategic\s+(?:Implications?|Considerations?)[\s:]*\n((?:[-•*]\s*.+\n?)+)",
        ]

        for pattern in key_section_patterns:
            matches = re.findall(pattern, brief_markdown, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                bullets = re.findall(r"[-•*]\s*(.+?)(?:\n|$)", match)
                takeaways.extend([b.strip() for b in bullets if len(b.strip()) > 20])

        # If no structured takeaways found, extract from summary section
        if not takeaways:
            if summary_match := re.search(
                r"(?:##?\s*)?(?:Executive\s+)?Summary[\s:]*\n(.+?)(?:\n##|\n\n\n|$)",
                brief_markdown,
                re.IGNORECASE | re.DOTALL,
            ):
                summary = summary_match.group(1)
                sentences = re.split(r"(?<=[.!?])\s+", summary)
                takeaways = [s.strip() for s in sentences[:3] if len(s.strip()) > 30]

        return takeaways[:5]  # Limit to 5 takeaways

    def _extract_city_examples(self, brief_markdown: str) -> List[Dict[str, str]]:
        """
        Extract examples of other cities, projects, or implementations from brief content.

        Returns list of dicts with 'city', 'project', and 'detail' keys.
        """

        examples = []

        # Common city/organization patterns
        city_patterns = [
            # "City of X has implemented/launched/deployed..."
            r"(?:City\s+of\s+|The\s+)?([\w\s]+?)(?:\s+has|\s+is|\s+launched|\s+implemented|\s+deployed|\s+piloted|\s+tested)\s+(.+?)(?:\.|,\s+(?:which|resulting|leading))",
            # "In X, they have..."
            r"In\s+([\w\s,]+?),\s+(?:they|the\s+city|officials|government)\s+(?:have|has)\s+(.+?)(?:\.|,)",
            # "X's program/initiative/project..."
            r"([\w\s]+?)'s\s+([\w\s]+?(?:program|initiative|project|pilot|system))\s+(.+?)(?:\.|,)",
            # "programs like X in Y"
            r"(?:programs?|initiatives?|projects?)\s+(?:like|such\s+as)\s+(.+?)\s+in\s+([\w\s]+?)(?:\.|,|$)",
        ]

        for pattern in city_patterns:
            matches = re.findall(pattern, brief_markdown, re.IGNORECASE)
            for match in matches:
                if len(match) >= 2:
                    city = match[0].strip() if match[0] else ""
                    detail = match[1].strip() if len(match) > 1 else ""
                    project = match[2].strip() if len(match) > 2 else ""

                    # Filter out generic terms
                    skip_terms = [
                        "the",
                        "this",
                        "that",
                        "these",
                        "those",
                        "austin",
                        "texas",
                    ]
                    if city.lower() not in skip_terms and len(city) > 2:
                        examples.append(
                            {
                                "city": city,
                                "project": project,
                                "detail": detail[:150] if detail else "",
                            }
                        )

        # Deduplicate by city name
        seen_cities = set()
        unique_examples = []
        for ex in examples:
            city_lower = ex["city"].lower()
            if city_lower not in seen_cities:
                seen_cities.add(city_lower)
                unique_examples.append(ex)

        return unique_examples[:4]  # Limit to 4 examples

    def _generate_portfolio_comparison_chart(
        self, briefs: List, dpi: int = CHART_DPI  # List of PortfolioBrief
    ) -> Optional[str]:
        """
        Generate a comparison chart showing all cards' scores.

        Creates a grouped bar chart comparing impact/relevance/velocity
        across all portfolio cards.
        """
        try:
            # Filter briefs with valid scores
            valid_briefs = [
                b
                for b in briefs
                if b.impact_score is not None or b.relevance_score is not None
            ]

            if not valid_briefs:
                return None

            fig, ax = plt.subplots(figsize=(10, 6))

            # Prepare data
            names = [
                f"{b.card_name[:25]}..." if len(b.card_name) > 25 else b.card_name
                for b in valid_briefs
            ]
            impacts = [b.impact_score or 0 for b in valid_briefs]
            relevances = [b.relevance_score or 0 for b in valid_briefs]
            velocities = [b.velocity_score or 0 for b in valid_briefs]

            x = np.arange(len(names))
            width = 0.25

            # Create bars
            ax.bar(
                x - width,
                impacts,
                width,
                label="Impact",
                color=COA_BRAND_COLORS["logo_blue"],
            )
            ax.bar(
                x,
                relevances,
                width,
                label="Relevance",
                color=COA_BRAND_COLORS["logo_green"],
            )
            ax.bar(
                x + width,
                velocities,
                width,
                label="Velocity",
                color=COA_BRAND_COLORS["dark_blue"],
            )

            # Customize chart
            ax.set_ylabel("Score (0-100)", fontsize=11)
            ax.set_title(
                "Portfolio Score Comparison",
                fontsize=14,
                fontweight="bold",
                color=COA_BRAND_COLORS["dark_blue"],
            )
            ax.set_xticks(x)
            ax.set_xticklabels(names, rotation=45, ha="right", fontsize=9)
            ax.legend(loc="upper right")
            ax.set_ylim(0, 110)

            # Style
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)
            ax.yaxis.grid(True, linestyle="--", alpha=0.3)

            plt.tight_layout()

            # Save
            temp_file = tempfile.NamedTemporaryFile(
                suffix=".png", delete=False, prefix="foresight_portfolio_comparison_"
            )
            plt.savefig(temp_file.name, dpi=dpi, bbox_inches="tight", facecolor="white")

            return temp_file.name

        except Exception as e:
            logger.error(f"Error generating portfolio comparison chart: {e}")
            return None
        finally:
            plt.close("all")

    def _generate_priority_matrix_chart(
        self,
        briefs: List,  # List of PortfolioBrief
        synthesis,  # PortfolioSynthesisData
        dpi: int = CHART_DPI,
    ) -> Optional[str]:
        """
        Generate a visual 2x2 priority matrix chart.

        Places cards in quadrants based on synthesis priority_matrix data.
        """
        try:
            matrix = synthesis.priority_matrix or {}
            urgent = set(matrix.get("high_impact_urgent", []))
            strategic = set(matrix.get("high_impact_strategic", []))
            monitor = set(matrix.get("monitor", []))

            fig, ax = plt.subplots(figsize=(10, 8))

            # Draw quadrant backgrounds
            ax.fill(
                [0, 50, 50, 0], [50, 50, 100, 100], color="#FEE2E2", alpha=0.5
            )  # Urgent - top left
            ax.fill(
                [50, 100, 100, 50], [50, 50, 100, 100], color="#FEF3C7", alpha=0.5
            )  # Strategic - top right
            ax.fill(
                [0, 50, 50, 0], [0, 0, 50, 50], color="#DBEAFE", alpha=0.5
            )  # Monitor - bottom left
            ax.fill(
                [50, 100, 100, 50], [0, 0, 50, 50], color="#D1FAE5", alpha=0.5
            )  # Low priority - bottom right

            # Quadrant labels
            ax.text(
                25,
                95,
                "🔴 URGENT ACTION",
                ha="center",
                va="top",
                fontsize=12,
                fontweight="bold",
                color="#DC2626",
            )
            ax.text(
                75,
                95,
                "🟡 STRATEGIC PLANNING",
                ha="center",
                va="top",
                fontsize=12,
                fontweight="bold",
                color="#D97706",
            )
            ax.text(
                25,
                5,
                "🔵 MONITOR",
                ha="center",
                va="bottom",
                fontsize=12,
                fontweight="bold",
                color="#2563EB",
            )
            ax.text(
                75,
                5,
                "🟢 EVALUATE",
                ha="center",
                va="bottom",
                fontsize=12,
                fontweight="bold",
                color="#059669",
            )

            # Place cards
            for i, brief in enumerate(briefs):
                name = brief.card_name
                # Determine quadrant based on synthesis
                if name in urgent:
                    x = 10 + (i % 3) * 12
                    y = 70 + (i // 3) * 8
                elif name in strategic:
                    x = 60 + (i % 3) * 12
                    y = 70 + (i // 3) * 8
                elif name in monitor:
                    x = 10 + (i % 3) * 12
                    y = 25 + (i // 3) * 8
                else:
                    x = 60 + (i % 3) * 12
                    y = 25 + (i // 3) * 8

                # Get pillar color
                pillar_def = PILLAR_DEFINITIONS.get(
                    brief.pillar_id.upper() if brief.pillar_id else "", {}
                )
                pillar_color = pillar_def.get("color", COA_BRAND_COLORS["logo_blue"])

                # Draw card marker
                ax.scatter(
                    x,
                    y,
                    s=200,
                    c=pillar_color,
                    edgecolors="white",
                    linewidth=2,
                    zorder=5,
                )

                # Truncate name
                short_name = f"{name[:18]}.." if len(name) > 18 else name
                ax.annotate(
                    short_name,
                    (x, y),
                    xytext=(0, -15),
                    textcoords="offset points",
                    ha="center",
                    fontsize=8,
                    color=COA_BRAND_COLORS["dark_blue"],
                )

            # Draw quadrant lines
            ax.axhline(y=50, color="gray", linewidth=2, linestyle="-", alpha=0.5)
            ax.axvline(x=50, color="gray", linewidth=2, linestyle="-", alpha=0.5)

            # Axis labels
            ax.set_xlabel(
                "← Lower Urgency          Higher Urgency →", fontsize=11, color="gray"
            )
            ax.set_ylabel(
                "← Lower Impact          Higher Impact →", fontsize=11, color="gray"
            )
            ax.set_title(
                "Strategic Priority Matrix",
                fontsize=14,
                fontweight="bold",
                color=COA_BRAND_COLORS["dark_blue"],
                pad=20,
            )

            ax.set_xlim(0, 100)
            ax.set_ylim(0, 100)
            ax.set_xticks([])
            ax.set_yticks([])

            for spine in ax.spines.values():
                spine.set_visible(False)

            plt.tight_layout()

            temp_file = tempfile.NamedTemporaryFile(
                suffix=".png", delete=False, prefix="foresight_priority_matrix_"
            )
            plt.savefig(temp_file.name, dpi=dpi, bbox_inches="tight", facecolor="white")

            return temp_file.name

        except Exception as e:
            logger.error(f"Error generating priority matrix chart: {e}")
            return None
        finally:
            plt.close("all")

    def _add_portfolio_dashboard_slide(
        self,
        prs: Presentation,
        briefs: List,  # List of PortfolioBrief
        comparison_chart_path: Optional[str],
        pillar_chart_path: Optional[str],
    ) -> None:
        """Add a visual dashboard slide with charts and key metrics."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_pptx_header(slide)
        self._add_pptx_footer(slide)

        # Title
        title_box = slide.shapes.add_textbox(
            PPTX_MARGIN, Inches(1.25), PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Portfolio Dashboard"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["primary"])

        # Key metrics row
        metrics_y = Inches(1.85)
        metric_width = Inches(2.2)
        metric_height = Inches(0.9)

        # Calculate metrics
        total_cards = len(briefs)
        avg_impact = sum(b.impact_score or 0 for b in briefs) // max(total_cards, 1)
        pillars_covered = len({b.pillar_id for b in briefs if b.pillar_id})
        horizons = {b.horizon for b in briefs if b.horizon}

        metrics = [
            (str(total_cards), "Strategic Trends"),
            (f"{avg_impact}/100", "Avg Impact Score") if avg_impact > 0 else None,
            (str(pillars_covered), "Pillars Covered"),
            (", ".join(sorted(horizons)) if horizons else "Mixed", "Time Horizons"),
        ]
        metrics = [m for m in metrics if m]  # Remove None

        for i, (value, label) in enumerate(metrics):
            x = PPTX_MARGIN + (i * (metric_width + Inches(0.15)))

            # Metric box
            metric_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, metrics_y, metric_width, metric_height
            )
            metric_box.fill.solid()
            metric_box.fill.fore_color.rgb = self._hex_to_rgb(
                COA_BRAND_COLORS["light_blue"]
            )
            metric_box.line.fill.background()

            # Value
            value_box = slide.shapes.add_textbox(
                x, metrics_y + Inches(0.1), metric_width, Inches(0.45)
            )
            value_frame = value_box.text_frame
            value_para = value_frame.paragraphs[0]
            value_para.text = value
            value_para.font.size = Pt(24)
            value_para.font.bold = True
            value_para.font.color.rgb = self._hex_to_rgb(COA_BRAND_COLORS["logo_blue"])
            value_para.alignment = PP_ALIGN.CENTER

            # Label
            label_box = slide.shapes.add_textbox(
                x, metrics_y + Inches(0.5), metric_width, Inches(0.35)
            )
            label_frame = label_box.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = label
            label_para.font.size = Pt(11)
            label_para.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])
            label_para.alignment = PP_ALIGN.CENTER

        # Charts row
        chart_y = Inches(3.0)
        chart_height = Inches(3.3)

        if comparison_chart_path:
            try:
                slide.shapes.add_picture(
                    comparison_chart_path,
                    PPTX_MARGIN,
                    chart_y,
                    width=Inches(4.5),
                    height=chart_height,
                )
            except Exception as e:
                logger.warning(f"Failed to add comparison chart: {e}")

        if pillar_chart_path:
            try:
                slide.shapes.add_picture(
                    pillar_chart_path,
                    Inches(5.0),
                    chart_y,
                    width=Inches(4.2),
                    height=chart_height,
                )
            except Exception as e:
                logger.warning(f"Failed to add pillar chart: {e}")

    def _add_card_deep_dive_slides(
        self,
        prs: Presentation,
        brief,  # PortfolioBrief
        index: int,
        chart_path: Optional[str] = None,
    ) -> None:
        """
        Add 2-3 slides for a single card with detailed insights.

        Slide 1: Overview with pillar, horizon, scores (if available)
        Slide 2: Key takeaways and city examples (if found)
        """
        pillar_def = PILLAR_DEFINITIONS.get(
            brief.pillar_id.upper() if brief.pillar_id else "", {}
        )
        pillar_name = pillar_def.get("name", brief.pillar_id or "Unknown")
        pillar_icon = pillar_def.get("icon", "🏛️")
        pillar_color = pillar_def.get("color", COA_BRAND_COLORS["logo_blue"])

        horizon_def = HORIZON_DEFINITIONS.get(
            brief.horizon.upper() if brief.horizon else "", {}
        )
        horizon_name = horizon_def.get("name", brief.horizon or "Unknown")

        stage_def = STAGE_DEFINITIONS.get(
            brief.stage_id.upper() if brief.stage_id else "", {}
        )
        stage_name = stage_def.get("name", brief.stage_id or "Unknown")

        # ===== SLIDE 1: Overview =====
        slide_layout = prs.slide_layouts[6]
        slide1 = prs.slides.add_slide(slide_layout)

        self._add_pptx_header(slide1)
        self._add_pptx_footer(slide1)

        # Title with index
        title_box = slide1.shapes.add_textbox(
            PPTX_MARGIN, Inches(1.25), PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        card_title = f"{index}. {brief.card_name}"
        title_para.text = card_title[:55] if len(card_title) > 55 else card_title
        title_para.font.size = Pt(26)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["primary"])

        # Classification badges row
        badges_y = Inches(1.9)
        badge_height = Inches(0.4)

        # Pillar badge
        pillar_badge = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            PPTX_MARGIN,
            badges_y,
            Inches(2.5),
            badge_height,
        )
        pillar_badge.fill.solid()
        pillar_badge.fill.fore_color.rgb = self._hex_to_rgb(pillar_color)
        pillar_badge.line.fill.background()

        pillar_text = slide1.shapes.add_textbox(
            PPTX_MARGIN, badges_y, Inches(2.5), badge_height
        )
        pf = pillar_text.text_frame
        pp = pf.paragraphs[0]
        pp.text = f"{pillar_icon} {pillar_name}"
        pp.font.size = Pt(14)
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(255, 255, 255)
        pp.alignment = PP_ALIGN.CENTER
        pf.paragraphs[0].space_before = Pt(8)

        # Horizon badge
        horizon_badge = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3.0),
            badges_y,
            Inches(2.0),
            badge_height,
        )
        horizon_badge.fill.solid()
        horizon_badge.fill.fore_color.rgb = self._hex_to_rgb(
            COA_BRAND_COLORS["dark_blue"]
        )
        horizon_badge.line.fill.background()

        horizon_text = slide1.shapes.add_textbox(
            Inches(3.0), badges_y, Inches(2.0), badge_height
        )
        hf = horizon_text.text_frame
        hp = hf.paragraphs[0]
        hp.text = f"⏱️ {horizon_name}"
        hp.font.size = Pt(14)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(255, 255, 255)
        hp.alignment = PP_ALIGN.CENTER
        hf.paragraphs[0].space_before = Pt(8)

        # Stage badge
        stage_badge = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.3),
            badges_y,
            Inches(2.2),
            badge_height,
        )
        stage_badge.fill.solid()
        stage_badge.fill.fore_color.rgb = self._hex_to_rgb(
            COA_BRAND_COLORS["logo_green"]
        )
        stage_badge.line.fill.background()

        stage_text = slide1.shapes.add_textbox(
            Inches(5.3), badges_y, Inches(2.2), badge_height
        )
        sf = stage_text.text_frame
        sp = sf.paragraphs[0]
        sp.text = f"📊 {stage_name}"
        sp.font.size = Pt(14)
        sp.font.bold = True
        sp.font.color.rgb = RGBColor(255, 255, 255)
        sp.alignment = PP_ALIGN.CENTER
        sf.paragraphs[0].space_before = Pt(8)

        # Summary content
        summary_y = Inches(2.5)
        summary_height = Inches(2.2)

        # If we have a chart, put it on the right
        if chart_path:
            summary_width = Inches(4.8)
            summary_box = slide1.shapes.add_textbox(
                PPTX_MARGIN, summary_y, summary_width, summary_height
            )

            # Add chart
            try:
                slide1.shapes.add_picture(
                    chart_path,
                    Inches(5.3),
                    summary_y,
                    width=Inches(4.0),
                    height=Inches(2.8),
                )
            except Exception as e:
                logger.warning(f"Failed to add score chart for {brief.card_name}: {e}")
        else:
            summary_width = PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN)
            summary_box = slide1.shapes.add_textbox(
                PPTX_MARGIN, summary_y, summary_width, summary_height
            )

        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True

        # Add scores if available
        scores_line = []
        if brief.impact_score and brief.impact_score > 0:
            scores_line.append(f"Impact: {brief.impact_score}/100")
        if brief.relevance_score and brief.relevance_score > 0:
            scores_line.append(f"Relevance: {brief.relevance_score}/100")
        if brief.velocity_score and brief.velocity_score > 0:
            scores_line.append(f"Velocity: {brief.velocity_score}/100")

        if scores_line:
            scores_para = summary_frame.paragraphs[0]
            scores_para.text = " | ".join(scores_line)
            scores_para.font.size = Pt(12)
            scores_para.font.bold = True
            scores_para.font.color.rgb = self._hex_to_rgb(COA_BRAND_COLORS["logo_blue"])
            scores_para.space_after = Pt(12)

            summary_para = summary_frame.add_paragraph()
        else:
            summary_para = summary_frame.paragraphs[0]

        summary_text = brief.brief_summary or "Executive summary not available."
        summary_para.text = (
            summary_text[:600] if len(summary_text) > 600 else summary_text
        )
        summary_para.font.size = Pt(14)
        summary_para.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])
        summary_para.space_before = Pt(6)

        # ===== SLIDE 2: Key Takeaways & Examples =====
        takeaways = self._extract_key_takeaways(brief.brief_content_markdown)
        city_examples = self._extract_city_examples(brief.brief_content_markdown)

        # Only add this slide if we have content
        if takeaways or city_examples:
            slide2 = prs.slides.add_slide(slide_layout)

            self._add_pptx_header(slide2)
            self._add_pptx_footer(slide2)

            # Title
            title_box2 = slide2.shapes.add_textbox(
                PPTX_MARGIN,
                Inches(1.25),
                PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN),
                Inches(0.5),
            )
            title_frame2 = title_box2.text_frame
            title_para2 = title_frame2.paragraphs[0]
            title_para2.text = f"{brief.card_name[:40]} - Key Insights"
            title_para2.font.size = Pt(24)
            title_para2.font.bold = True
            title_para2.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["primary"])

            content_y = Inches(1.85)

        # Key Takeaways section
        if takeaways:
            takeaways_box = slide2.shapes.add_textbox(
                PPTX_MARGIN, content_y, Inches(4.5), Inches(4.0)
            )
            tf = takeaways_box.text_frame
            tf.word_wrap = True

            # Section header
            header_para = tf.paragraphs[0]
            header_para.text = "📌 Key Takeaways"
            header_para.font.size = Pt(16)
            header_para.font.bold = True
            header_para.font.color.rgb = self._hex_to_rgb(COA_BRAND_COLORS["logo_blue"])
            header_para.space_after = Pt(8)

            for takeaway in takeaways:
                bullet_para = tf.add_paragraph()
                bullet_text = takeaway[:200] if len(takeaway) > 200 else takeaway
                bullet_para.text = f"• {bullet_text}"
                bullet_para.font.size = Pt(12)
                bullet_para.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])
                bullet_para.space_before = Pt(6)

        # City Examples section
        if city_examples:
            examples_x = Inches(5.0) if takeaways else PPTX_MARGIN
            examples_width = (
                Inches(4.2) if takeaways else PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN)
            )

            examples_box = slide2.shapes.add_textbox(
                examples_x, content_y, examples_width, Inches(4.0)
            )
            ef = examples_box.text_frame
            ef.word_wrap = True

            # Section header
            ex_header = ef.paragraphs[0]
            ex_header.text = "🌆 Examples from Other Cities"
            ex_header.font.size = Pt(16)
            ex_header.font.bold = True
            ex_header.font.color.rgb = self._hex_to_rgb(COA_BRAND_COLORS["logo_green"])
            ex_header.space_after = Pt(8)

            for example in city_examples:
                city_para = ef.add_paragraph()
                city_text = f"• {example['city']}"
                if example.get("detail"):
                    city_text += f": {example['detail']}"
                city_para.text = city_text[:180] if len(city_text) > 180 else city_text
                city_para.font.size = Pt(12)
                city_para.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])
                city_para.space_before = Pt(6)

    async def generate_portfolio_pptx_local(
        self,
        workstream_name: str,
        briefs: List,  # List of PortfolioBrief
        synthesis,  # PortfolioSynthesisData
    ) -> str:
        """
        Generate a local portfolio PPTX presentation.

        Creates a professional multi-card portfolio deck with:
        - Title slide with branding
        - Portfolio dashboard with metrics and charts
        - Executive overview synthesis
        - Visual priority matrix (2x2)
        - Per-card deep dives (2-3 slides each) with:
          - Classification badges and scores (when available)
          - Key takeaways extracted from brief
          - City/project examples cited in research
        - Cross-cutting themes
        - Recommended actions
        - AI disclosure

        Args:
            workstream_name: Name of the workstream for title
            briefs: List of PortfolioBrief objects
            synthesis: PortfolioSynthesisData with AI-generated overview

        Returns:
            Path to the generated PPTX file
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from datetime import datetime, timezone
        import tempfile

        prs = Presentation()
        prs.slide_width = PPTX_SLIDE_WIDTH
        prs.slide_height = PPTX_SLIDE_HEIGHT

        temp_files_to_cleanup = []

        try:
            # Get pillar icons for title
            pillar_icons = []
            pillar_counts = {}
            for brief in briefs:
                pillar_def = PILLAR_DEFINITIONS.get(
                    brief.pillar_id.upper() if brief.pillar_id else "", {}
                )
                icon = pillar_def.get("icon", "🏛️")
                pillar_name = pillar_def.get("name", brief.pillar_id or "Other")
                if icon not in pillar_icons:
                    pillar_icons.append(icon)
                pillar_counts[pillar_name] = pillar_counts.get(pillar_name, 0) + 1

            # ===== 1. TITLE SLIDE =====
            title_subtitle = f"{' '.join(pillar_icons)} | {len(briefs)} Strategic Trends\n{datetime.now(timezone.utc).strftime('%B %Y')}"
            self._add_title_slide(prs, workstream_name, title_subtitle)

            # ===== 2. PORTFOLIO DASHBOARD =====
            # Generate charts
            comparison_chart_path = self._generate_portfolio_comparison_chart(briefs)
            if comparison_chart_path:
                temp_files_to_cleanup.append(comparison_chart_path)

            pillar_chart_path = None
            if pillar_counts:
                pillar_chart_path = self.generate_pillar_distribution_chart(
                    pillar_counts, "Distribution by Pillar"
                )
                if pillar_chart_path:
                    temp_files_to_cleanup.append(pillar_chart_path)

            self._add_portfolio_dashboard_slide(
                prs, briefs, comparison_chart_path, pillar_chart_path
            )

            # ===== 3. WHY THIS MATTERS NOW =====
            urgency = (
                getattr(synthesis, "urgency_statement", "")
                or f"These {len(briefs)} trends represent critical opportunities and challenges. Early action positions Austin as a leader; delay risks falling behind peer cities."
            )
            urgency_content = f"{urgency}\n\n**The Window of Opportunity**\n\nCities that move first on emerging trends gain competitive advantage in talent attraction, federal funding, and citizen satisfaction."
            self._add_smart_content_slide(
                prs,
                title="Why This Matters Now",
                content=urgency_content,
                max_chars=1200,
            )

            # ===== 4. EXECUTIVE OVERVIEW =====
            overview_content = (
                synthesis.executive_overview or "Portfolio synthesis in progress..."
            )
            self._add_smart_content_slide(
                prs,
                title="Executive Overview",
                content=overview_content,
                max_chars=1400,
            )

            if matrix_chart_path := self._generate_priority_matrix_chart(
                briefs, synthesis
            ):
                temp_files_to_cleanup.append(matrix_chart_path)

                # Add matrix slide
                slide_layout = prs.slide_layouts[6]
                matrix_slide = prs.slides.add_slide(slide_layout)
                self._add_pptx_header(matrix_slide)
                self._add_pptx_footer(matrix_slide)

                title_box = matrix_slide.shapes.add_textbox(
                    PPTX_MARGIN,
                    Inches(1.25),
                    PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN),
                    Inches(0.5),
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = "Strategic Priority Matrix"
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                title_para.font.color.rgb = self._hex_to_rgb(
                    FORESIGHT_COLORS["primary"]
                )

                try:
                    matrix_slide.shapes.add_picture(
                        matrix_chart_path,
                        Inches(0.8),
                        Inches(1.9),
                        width=Inches(8.4),
                        height=Inches(5.0),
                    )
                except Exception as e:
                    logger.warning(f"Failed to add priority matrix chart: {e}")
            else:
                # Fallback to text-based priority slide
                matrix = synthesis.priority_matrix or {}
                urgent = matrix.get("high_impact_urgent", [])
                strategic = matrix.get("high_impact_strategic", [])
                monitor = matrix.get("monitor", [])

                priority_content = "🔴 **High Impact - Urgent Action**\n" + (
                    "\n".join(f"• {item}" for item in urgent)
                    if urgent
                    else "• None identified"
                )
                priority_content += "\n\n🟡 **High Impact - Strategic Planning**\n"
                priority_content += (
                    "\n".join(f"• {item}" for item in strategic)
                    if strategic
                    else "• None identified"
                )
                priority_content += "\n\n🟢 **Monitor & Evaluate**\n"
                priority_content += (
                    "\n".join(f"• {item}" for item in monitor)
                    if monitor
                    else "• None identified"
                )

                self._add_smart_content_slide(
                    prs,
                    title="Strategic Priorities",
                    content=priority_content,
                    max_chars=1500,
                )

            # ===== 6. IMPLEMENTATION GUIDANCE =====
            impl = getattr(synthesis, "implementation_guidance", {}) or {}
            impl_lines = []
            if impl.get("pilot_now"):
                impl_lines.append(
                    f"🚀 **Ready to Pilot**: {', '.join(impl['pilot_now'])}"
                )
            if impl.get("investigate_further"):
                impl_lines.append(
                    f"🔍 **Investigate Further**: {', '.join(impl['investigate_further'])}"
                )
            if impl.get("meet_with_vendors"):
                impl_lines.append(
                    f"🤝 **Meet with Vendors**: {', '.join(impl['meet_with_vendors'])}"
                )
            if impl.get("policy_review"):
                impl_lines.append(
                    f"📋 **Policy Review Needed**: {', '.join(impl['policy_review'])}"
                )
            if impl.get("staff_training"):
                impl_lines.append(
                    f"👥 **Staff Training Focus**: {', '.join(impl['staff_training'])}"
                )
            if impl.get("budget_planning"):
                impl_lines.append(
                    f"💰 **Budget Planning**: {', '.join(impl['budget_planning'])}"
                )

            if impl_lines:
                impl_content = "What should Austin DO with each trend?\n\n" + "\n".join(
                    impl_lines
                )
                self._add_smart_content_slide(
                    prs,
                    title="Implementation Guidance",
                    content=impl_content,
                    max_chars=1400,
                )

            # ===== 7. PER-CARD DEEP DIVES =====
            for i, brief in enumerate(briefs, 1):
                # Generate score chart for this card if scores exist
                card_chart_path = None
                if (
                    (brief.impact_score and brief.impact_score > 0)
                    or (brief.relevance_score and brief.relevance_score > 0)
                    or (brief.velocity_score and brief.velocity_score > 0)
                ):
                    # Create a simple CardExportData-like object for chart generation
                    scores = {
                        "Impact": brief.impact_score or 0,
                        "Relevance": brief.relevance_score or 0,
                        "Velocity": brief.velocity_score or 0,
                    }
                    if valid_scores := {k: v for k, v in scores.items() if v > 0}:
                        card_chart_path = self._generate_radar_chart(
                            valid_scores, brief.card_name, CHART_DPI
                        )
                        if card_chart_path:
                            temp_files_to_cleanup.append(card_chart_path)

                self._add_card_deep_dive_slides(prs, brief, i, card_chart_path)

            # ===== 8. CROSS-CUTTING THEMES =====
            themes_content = "**Common Patterns Across Trends**\n"
            themes_content += (
                "\n".join(f"• {theme}" for theme in (synthesis.key_themes or []))
                or "• Analysis in progress"
            )
            themes_content += "\n\n**Strategic Connections**\n"
            themes_content += (
                "\n".join(
                    f"• {insight}"
                    for insight in (synthesis.cross_cutting_insights or [])
                )
                or "• Analysis in progress"
            )

            self._add_smart_content_slide(
                prs,
                title="Cross-Cutting Themes",
                content=themes_content,
                max_chars=1400,
            )

            if ninety_day := getattr(synthesis, "ninety_day_actions", []) or []:
                actions_content = "What Austin should do in the next 90 days:\n\n"
                for action in ninety_day[:5]:
                    action_text = action.get("action", "")
                    owner = action.get("owner", "TBD")
                    by_when = action.get("by_when", "90 days")
                    metric = action.get("success_metric", "")
                    actions_content += f"✓ **{action_text}**\n"
                    actions_content += f"   Owner: {owner} | By: {by_when}"
                    if metric:
                        actions_content += f"\n   Success: {metric}"
                    actions_content += "\n\n"
            else:
                # Fall back to recommended_actions
                actions_content = ""
                for action in (synthesis.recommended_actions or [])[:6]:
                    action_text = action.get("action", "")
                    owner = action.get("owner", "TBD")
                    timeline = action.get("timeline", "TBD")
                    related_cards = action.get("cards", [])

                    actions_content += f"✓ **{action_text}**\n"
                    actions_content += f"   Owner: {owner} | Timeline: {timeline}"
                    if related_cards:
                        actions_content += f" | Related: {', '.join(related_cards[:2])}"
                    actions_content += "\n\n"

                if not actions_content:
                    actions_content = (
                        "Action plan to be developed based on leadership priorities."
                    )

            self._add_smart_content_slide(
                prs, title="90-Day Action Plan", content=actions_content, max_chars=1400
            )

            # ===== 10. RISKS & OPPORTUNITIES =====
            risk_text = (
                getattr(synthesis, "risk_summary", "")
                or "Delayed action on these trends could result in Austin falling behind peer cities, missing federal funding windows, and losing competitive advantage."
            )
            opp_text = (
                getattr(synthesis, "opportunity_summary", "")
                or "Early action positions Austin as a national leader, attracts innovation investment, and delivers improved services to residents."
            )

            risk_opp_content = f"⚠️ **If Austin Doesn't Act**\n{risk_text}\n\n✨ **If Austin Leads**\n{opp_text}"
            self._add_smart_content_slide(
                prs,
                title="Risks & Opportunities",
                content=risk_opp_content,
                max_chars=1400,
            )

            # ===== 11. AI DISCLOSURE =====
            self._add_ai_disclosure_slide(prs)

            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(
                suffix=".pptx", delete=False, prefix="foresight_portfolio_local_"
            )
            prs.save(temp_file.name)
            temp_file.close()

            logger.info(
                f"Generated enhanced local portfolio PPTX: {len(briefs)} cards, {len(prs.slides)} slides"
            )
            return temp_file.name

        finally:
            # Clean up temp chart files
            self.cleanup_temp_files(temp_files_to_cleanup)

    async def generate_portfolio_pdf(
        self,
        workstream_name: str,
        briefs: List,  # List of PortfolioBrief
        synthesis,  # PortfolioSynthesisData
    ) -> str:
        """
        Generate a detailed portfolio PDF document.

        PDF version includes MORE detail than PPTX since readers can
        absorb more information. Same structure but expanded content.

        Args:
            workstream_name: Name of the workstream for title
            briefs: List of PortfolioBrief objects
            synthesis: PortfolioSynthesisData with AI-generated overview

        Returns:
            Path to the generated PDF file
        """
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.colors import HexColor
        from reportlab.platypus import (
            Paragraph,
            Spacer,
            PageBreak,
        )
        from reportlab.lib.units import inch
        from datetime import datetime, timezone
        import tempfile

        # Create temp file
        temp_file = tempfile.NamedTemporaryFile(
            suffix=".pdf", delete=False, prefix="foresight_portfolio_"
        )

        doc = SimpleDocTemplate(
            temp_file.name,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72,
        )

        styles = getSampleStyleSheet()

        # Custom styles
        title_style = ParagraphStyle(
            "PortfolioTitle",
            parent=styles["Heading1"],
            fontSize=24,
            textColor=HexColor(COA_BRAND_COLORS["logo_blue"]),
            spaceAfter=12,
            alignment=1,  # Center
        )

        subtitle_style = ParagraphStyle(
            "PortfolioSubtitle",
            parent=styles["Normal"],
            fontSize=14,
            textColor=HexColor(COA_BRAND_COLORS["dark_gray"]),
            spaceAfter=24,
            alignment=1,
        )

        section_style = ParagraphStyle(
            "SectionTitle",
            parent=styles["Heading2"],
            fontSize=16,
            textColor=HexColor(COA_BRAND_COLORS["logo_blue"]),
            spaceBefore=18,
            spaceAfter=12,
        )

        card_title_style = ParagraphStyle(
            "CardTitle",
            parent=styles["Heading3"],
            fontSize=14,
            textColor=HexColor(COA_BRAND_COLORS["dark_blue"]),
            spaceBefore=12,
            spaceAfter=6,
        )

        body_style = ParagraphStyle(
            "BodyText",
            parent=styles["Normal"],
            fontSize=11,
            textColor=HexColor(COA_BRAND_COLORS["dark_gray"]),
            spaceAfter=8,
            leading=14,
        )

        elements = [Spacer(1, inch * 2)]

        elements.append(Paragraph(workstream_name, title_style))
        elements.append(Paragraph("Strategic Intelligence Portfolio", subtitle_style))
        elements.append(
            Paragraph(
                f"{len(briefs)} Strategic Trends | {datetime.now(timezone.utc).strftime('%B %Y')}",
                subtitle_style,
            )
        )
        elements.append(
            Paragraph("City of Austin | FORESIGHT Platform", subtitle_style)
        )
        elements.append(PageBreak())

        # Executive Overview
        elements.append(Paragraph("Executive Overview", section_style))
        overview_text = (
            synthesis.executive_overview or "Portfolio analysis in progress."
        )
        # Split into paragraphs for better formatting
        for para in overview_text.split("\n\n"):
            if para.strip():
                elements.append(Paragraph(para.strip(), body_style))
        elements.append(Spacer(1, 12))

        # Key Themes
        elements.append(Paragraph("Key Themes", section_style))
        elements.extend(
            Paragraph(f"• {theme}", body_style)
            for theme in (synthesis.key_themes or [])
        )
        elements.append(Spacer(1, 12))

        # Strategic Priorities
        elements.append(Paragraph("Strategic Priorities", section_style))
        matrix = synthesis.priority_matrix or {}

        if urgent := matrix.get("high_impact_urgent", []):
            elements.append(
                Paragraph("<b>High Impact - Urgent Action:</b>", body_style)
            )
            for item in urgent:
                elements.append(Paragraph(f"  • {item}", body_style))

        if strategic := matrix.get("high_impact_strategic", []):
            elements.append(
                Paragraph("<b>High Impact - Strategic Planning:</b>", body_style)
            )
            for item in strategic:
                elements.append(Paragraph(f"  • {item}", body_style))

        if monitor := matrix.get("monitor", []):
            elements.append(Paragraph("<b>Monitor & Evaluate:</b>", body_style))
            for item in monitor:
                elements.append(Paragraph(f"  • {item}", body_style))

        elements.append(PageBreak())

        # Per-card detailed sections (PDF gets FULL content)
        elements.append(Paragraph("Trend Analysis", section_style))

        for i, brief in enumerate(briefs, 1):
            pillar_def = PILLAR_DEFINITIONS.get(
                brief.pillar_id.upper() if brief.pillar_id else "", {}
            )
            pillar_name = pillar_def.get("name", brief.pillar_id or "Unknown")
            horizon_name = brief.horizon or "H2"

            elements.append(Paragraph(f"{i}. {brief.card_name}", card_title_style))
            elements.append(
                Paragraph(
                    f"<b>Pillar:</b> {pillar_name} | <b>Horizon:</b> {horizon_name} | "
                    f"<b>Impact:</b> {brief.impact_score}/100 | <b>Relevance:</b> {brief.relevance_score}/100",
                    body_style,
                )
            )

            # Summary
            if brief.brief_summary:
                elements.append(
                    Paragraph(f"<b>Summary:</b> {brief.brief_summary}", body_style)
                )

            # Full brief content (PDF gets expanded detail)
            if brief.brief_content_markdown:
                # Clean up markdown for PDF
                content = brief.brief_content_markdown
                # Remove markdown headers, keep content
                import re

                content = re.sub(r"^#+\s+", "", content, flags=re.MULTILINE)
                content = re.sub(r"\*\*([^*]+)\*\*", r"<b>\1</b>", content)
                content = re.sub(r"\*([^*]+)\*", r"<i>\1</i>", content)

                # Truncate if very long
                if len(content) > 2000:
                    content = f"{content[:2000]}..."

                for para in content.split("\n\n")[:5]:  # First 5 paragraphs
                    if para.strip():
                        elements.append(Paragraph(para.strip(), body_style))

            elements.append(Spacer(1, 12))

        elements.append(PageBreak())

        # Cross-Cutting Insights
        elements.append(Paragraph("Cross-Cutting Insights", section_style))
        for insight in synthesis.cross_cutting_insights or []:
            elements.append(Paragraph(f"• {insight}", body_style))
        elements.append(Spacer(1, 12))

        # Recommended Actions
        elements.append(Paragraph("Recommended Actions", section_style))
        for action in synthesis.recommended_actions or []:
            action_text = action.get("action", "")
            owner = action.get("owner", "TBD")
            timeline = action.get("timeline", "TBD")
            cards = ", ".join(action.get("cards", []))

            elements.append(Paragraph(f"<b>{action_text}</b>", body_style))
            elements.append(
                Paragraph(f"  Owner: {owner} | Timeline: {timeline}", body_style)
            )
            if cards:
                elements.append(Paragraph(f"  Related Trends: {cards}", body_style))
            elements.append(Spacer(1, 6))

        # AI Disclosure
        elements.append(PageBreak())
        elements.append(Paragraph("About This Portfolio", section_style))
        disclosure = f"""This strategic intelligence portfolio was generated using the FORESIGHT platform,
        powered by advanced AI technologies including OpenAI {get_chat_deployment()}, GPT Researcher,
        SearXNG, Serper, trafilatura, and Gamma.app. The City of Austin is committed to transparent
        and responsible use of AI technology in public service. All AI-generated content is reviewed
        for accuracy and relevance."""
        elements.append(Paragraph(disclosure, body_style))

        # Build PDF
        doc.build(elements)
        temp_file.close()

        logger.info(f"Generated portfolio PDF: {len(briefs)} cards")
        return temp_file.name
