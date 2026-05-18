"""Portfolio-specific PPTX slide builders.

Builds the dashboard slide (metrics + charts) and per-card deep-dive slides
(overview + key takeaways/city examples).
"""

import logging
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ...gamma_service import (
    HORIZON_DEFINITIONS,
    PILLAR_DEFINITIONS,
    STAGE_DEFINITIONS,
)
from ..branding import COA_BRAND_COLORS, FORESIGHT_COLORS
from ..pptx import (
    PPTX_MARGIN,
    PPTX_SLIDE_WIDTH,
    add_pptx_footer,
    add_pptx_header,
    hex_to_rgb,
)
from .extractors import extract_city_examples, extract_key_takeaways

logger = logging.getLogger(__name__)


def add_portfolio_dashboard_slide(
    prs: Presentation,
    briefs: List,  # List of PortfolioBrief
    comparison_chart_path: Optional[str],
    pillar_chart_path: Optional[str],
) -> None:
    """Add a visual dashboard slide with key metrics and two charts."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_pptx_header(slide)
    add_pptx_footer(slide)

    title_box = slide.shapes.add_textbox(
        PPTX_MARGIN, Inches(1.25), PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Portfolio Dashboard"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["primary"])

    metrics_y = Inches(1.85)
    metric_width = Inches(2.2)
    metric_height = Inches(0.9)

    total_cards = len(briefs)
    avg_impact = sum(b.impact_score or 0 for b in briefs) // max(total_cards, 1)
    pillars_covered = len({b.pillar_id for b in briefs if b.pillar_id})
    horizons = {b.horizon for b in briefs if b.horizon}

    metrics = [
        (str(total_cards), "Strategic Trends"),
        (f"{avg_impact}/100", "Avg Impact Score") if avg_impact > 0 else None,
        (str(pillars_covered), "Pillars Covered"),
        (", ".join(sorted(horizons)) if horizons else "Mixed", "Time Horizons"),
    ]
    metrics = [m for m in metrics if m]

    for i, (value, label) in enumerate(metrics):
        x = PPTX_MARGIN + (i * (metric_width + Inches(0.15)))

        metric_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, metrics_y, metric_width, metric_height
        )
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = hex_to_rgb(COA_BRAND_COLORS["light_blue"])
        metric_box.line.fill.background()

        value_box = slide.shapes.add_textbox(
            x, metrics_y + Inches(0.1), metric_width, Inches(0.45)
        )
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = value
        value_para.font.size = Pt(24)
        value_para.font.bold = True
        value_para.font.color.rgb = hex_to_rgb(COA_BRAND_COLORS["logo_blue"])
        value_para.alignment = PP_ALIGN.CENTER

        label_box = slide.shapes.add_textbox(
            x, metrics_y + Inches(0.5), metric_width, Inches(0.35)
        )
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = label
        label_para.font.size = Pt(11)
        label_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["dark"])
        label_para.alignment = PP_ALIGN.CENTER

    chart_y = Inches(3.0)
    chart_height = Inches(3.3)

    if comparison_chart_path:
        try:
            slide.shapes.add_picture(
                comparison_chart_path,
                PPTX_MARGIN,
                chart_y,
                width=Inches(4.5),
                height=chart_height,
            )
        except Exception as e:
            logger.warning(f"Failed to add comparison chart: {e}")

    if pillar_chart_path:
        try:
            slide.shapes.add_picture(
                pillar_chart_path,
                Inches(5.0),
                chart_y,
                width=Inches(4.2),
                height=chart_height,
            )
        except Exception as e:
            logger.warning(f"Failed to add pillar chart: {e}")


def add_card_deep_dive_slides(
    prs: Presentation,
    brief,  # PortfolioBrief
    index: int,
    chart_path: Optional[str] = None,
) -> None:
    """Add 2-3 slides for a single portfolio card.

    Slide 1: Overview with pillar/horizon/stage badges, scores, summary.
    Slide 2 (if content present): Key takeaways + examples from other cities.
    """
    pillar_def = PILLAR_DEFINITIONS.get(
        brief.pillar_id.upper() if brief.pillar_id else "", {}
    )
    pillar_name = pillar_def.get("name", brief.pillar_id or "Unknown")
    pillar_icon = pillar_def.get("icon", "🏛️")
    pillar_color = pillar_def.get("color", COA_BRAND_COLORS["logo_blue"])

    horizon_def = HORIZON_DEFINITIONS.get(
        brief.horizon.upper() if brief.horizon else "", {}
    )
    horizon_name = horizon_def.get("name", brief.horizon or "Unknown")

    stage_def = STAGE_DEFINITIONS.get(
        brief.stage_id.upper() if brief.stage_id else "", {}
    )
    stage_name = stage_def.get("name", brief.stage_id or "Unknown")

    # ===== SLIDE 1: Overview =====
    slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(slide_layout)

    add_pptx_header(slide1)
    add_pptx_footer(slide1)

    title_box = slide1.shapes.add_textbox(
        PPTX_MARGIN, Inches(1.25), PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    card_title = f"{index}. {brief.card_name}"
    title_para.text = card_title[:55] if len(card_title) > 55 else card_title
    title_para.font.size = Pt(26)
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["primary"])

    badges_y = Inches(1.9)
    badge_height = Inches(0.4)

    # Pillar badge
    pillar_badge = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        PPTX_MARGIN,
        badges_y,
        Inches(2.5),
        badge_height,
    )
    pillar_badge.fill.solid()
    pillar_badge.fill.fore_color.rgb = hex_to_rgb(pillar_color)
    pillar_badge.line.fill.background()

    pillar_text = slide1.shapes.add_textbox(
        PPTX_MARGIN, badges_y, Inches(2.5), badge_height
    )
    pf = pillar_text.text_frame
    pp = pf.paragraphs[0]
    pp.text = f"{pillar_icon} {pillar_name}"
    pp.font.size = Pt(14)
    pp.font.bold = True
    pp.font.color.rgb = RGBColor(255, 255, 255)
    pp.alignment = PP_ALIGN.CENTER
    pf.paragraphs[0].space_before = Pt(8)

    # Horizon badge
    horizon_badge = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.0),
        badges_y,
        Inches(2.0),
        badge_height,
    )
    horizon_badge.fill.solid()
    horizon_badge.fill.fore_color.rgb = hex_to_rgb(COA_BRAND_COLORS["dark_blue"])
    horizon_badge.line.fill.background()

    horizon_text = slide1.shapes.add_textbox(
        Inches(3.0), badges_y, Inches(2.0), badge_height
    )
    hf = horizon_text.text_frame
    hp = hf.paragraphs[0]
    hp.text = f"⏱️ {horizon_name}"
    hp.font.size = Pt(14)
    hp.font.bold = True
    hp.font.color.rgb = RGBColor(255, 255, 255)
    hp.alignment = PP_ALIGN.CENTER
    hf.paragraphs[0].space_before = Pt(8)

    # Stage badge
    stage_badge = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.3),
        badges_y,
        Inches(2.2),
        badge_height,
    )
    stage_badge.fill.solid()
    stage_badge.fill.fore_color.rgb = hex_to_rgb(COA_BRAND_COLORS["logo_green"])
    stage_badge.line.fill.background()

    stage_text = slide1.shapes.add_textbox(
        Inches(5.3), badges_y, Inches(2.2), badge_height
    )
    sf = stage_text.text_frame
    sp = sf.paragraphs[0]
    sp.text = f"📊 {stage_name}"
    sp.font.size = Pt(14)
    sp.font.bold = True
    sp.font.color.rgb = RGBColor(255, 255, 255)
    sp.alignment = PP_ALIGN.CENTER
    sf.paragraphs[0].space_before = Pt(8)

    # Summary content
    summary_y = Inches(2.5)
    summary_height = Inches(2.2)

    if chart_path:
        summary_width = Inches(4.8)
        summary_box = slide1.shapes.add_textbox(
            PPTX_MARGIN, summary_y, summary_width, summary_height
        )

        try:
            slide1.shapes.add_picture(
                chart_path,
                Inches(5.3),
                summary_y,
                width=Inches(4.0),
                height=Inches(2.8),
            )
        except Exception as e:
            logger.warning(f"Failed to add score chart for {brief.card_name}: {e}")
    else:
        summary_width = PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN)
        summary_box = slide1.shapes.add_textbox(
            PPTX_MARGIN, summary_y, summary_width, summary_height
        )

    summary_frame = summary_box.text_frame
    summary_frame.word_wrap = True

    scores_line = []
    if brief.impact_score and brief.impact_score > 0:
        scores_line.append(f"Impact: {brief.impact_score}/100")
    if brief.relevance_score and brief.relevance_score > 0:
        scores_line.append(f"Relevance: {brief.relevance_score}/100")
    if brief.velocity_score and brief.velocity_score > 0:
        scores_line.append(f"Velocity: {brief.velocity_score}/100")

    if scores_line:
        scores_para = summary_frame.paragraphs[0]
        scores_para.text = " | ".join(scores_line)
        scores_para.font.size = Pt(12)
        scores_para.font.bold = True
        scores_para.font.color.rgb = hex_to_rgb(COA_BRAND_COLORS["logo_blue"])
        scores_para.space_after = Pt(12)

        summary_para = summary_frame.add_paragraph()
    else:
        summary_para = summary_frame.paragraphs[0]

    summary_text = brief.brief_summary or "Executive summary not available."
    summary_para.text = (
        summary_text[:600] if len(summary_text) > 600 else summary_text
    )
    summary_para.font.size = Pt(14)
    summary_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["dark"])
    summary_para.space_before = Pt(6)

    # ===== SLIDE 2: Key Takeaways & Examples =====
    takeaways = extract_key_takeaways(brief.brief_content_markdown)
    city_examples = extract_city_examples(brief.brief_content_markdown)

    if not (takeaways or city_examples):
        return

    slide2 = prs.slides.add_slide(slide_layout)

    add_pptx_header(slide2)
    add_pptx_footer(slide2)

    title_box2 = slide2.shapes.add_textbox(
        PPTX_MARGIN,
        Inches(1.25),
        PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN),
        Inches(0.5),
    )
    title_frame2 = title_box2.text_frame
    title_para2 = title_frame2.paragraphs[0]
    title_para2.text = f"{brief.card_name[:40]} - Key Insights"
    title_para2.font.size = Pt(24)
    title_para2.font.bold = True
    title_para2.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["primary"])

    content_y = Inches(1.85)

    if takeaways:
        takeaways_box = slide2.shapes.add_textbox(
            PPTX_MARGIN, content_y, Inches(4.5), Inches(4.0)
        )
        tf = takeaways_box.text_frame
        tf.word_wrap = True

        header_para = tf.paragraphs[0]
        header_para.text = "📌 Key Takeaways"
        header_para.font.size = Pt(16)
        header_para.font.bold = True
        header_para.font.color.rgb = hex_to_rgb(COA_BRAND_COLORS["logo_blue"])
        header_para.space_after = Pt(8)

        for takeaway in takeaways:
            bullet_para = tf.add_paragraph()
            bullet_text = takeaway[:200] if len(takeaway) > 200 else takeaway
            bullet_para.text = f"• {bullet_text}"
            bullet_para.font.size = Pt(12)
            bullet_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["dark"])
            bullet_para.space_before = Pt(6)

    if city_examples:
        examples_x = Inches(5.0) if takeaways else PPTX_MARGIN
        examples_width = (
            Inches(4.2) if takeaways else PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN)
        )

        examples_box = slide2.shapes.add_textbox(
            examples_x, content_y, examples_width, Inches(4.0)
        )
        ef = examples_box.text_frame
        ef.word_wrap = True

        ex_header = ef.paragraphs[0]
        ex_header.text = "🌆 Examples from Other Cities"
        ex_header.font.size = Pt(16)
        ex_header.font.bold = True
        ex_header.font.color.rgb = hex_to_rgb(COA_BRAND_COLORS["logo_green"])
        ex_header.space_after = Pt(8)

        for example in city_examples:
            city_para = ef.add_paragraph()
            city_text = f"• {example['city']}"
            if example.get("detail"):
                city_text += f": {example['detail']}"
            city_para.text = city_text[:180] if len(city_text) > 180 else city_text
            city_para.font.size = Pt(12)
            city_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["dark"])
            city_para.space_before = Pt(6)
