"""Portfolio PPTX generator (local python-pptx).

Builds a full multi-card portfolio deck: title slide, dashboard with metrics +
charts, executive overview, priority matrix, per-card deep dives, themes,
actions, risks/opportunities, and AI disclosure.
"""

import logging
import tempfile
from datetime import datetime, timezone
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt

from ...gamma_service import PILLAR_DEFINITIONS
from ..briefs import add_ai_disclosure_slide, add_smart_content_slide
from ..branding import FORESIGHT_COLORS
from ..charts import (
    CHART_DPI,
    generate_pillar_distribution_chart,
    generate_radar_chart,
)
from ..pptx import (
    PPTX_MARGIN,
    PPTX_SLIDE_HEIGHT,
    PPTX_SLIDE_WIDTH,
    add_pptx_footer,
    add_pptx_header,
    add_title_slide,
    hex_to_rgb,
)
from ..utils import cleanup_temp_files
from .charts import (
    generate_portfolio_comparison_chart,
    generate_priority_matrix_chart,
)
from .slides import add_card_deep_dive_slides, add_portfolio_dashboard_slide

logger = logging.getLogger(__name__)


async def generate_portfolio_pptx_local(
    workstream_name: str,
    briefs: List,  # List of PortfolioBrief
    synthesis,  # PortfolioSynthesisData
) -> str:
    """Generate a local portfolio PPTX deck.

    Produces a professional multi-card portfolio with title, dashboard,
    executive overview, priority matrix, per-card deep dives, cross-cutting
    themes, recommended actions, risk/opportunity summary, and AI disclosure.

    Returns the path to the generated PPTX file.
    """
    prs = Presentation()
    prs.slide_width = PPTX_SLIDE_WIDTH
    prs.slide_height = PPTX_SLIDE_HEIGHT

    temp_files_to_cleanup = []

    try:
        pillar_icons = []
        pillar_counts = {}
        for brief in briefs:
            pillar_def = PILLAR_DEFINITIONS.get(
                brief.pillar_id.upper() if brief.pillar_id else "", {}
            )
            icon = pillar_def.get("icon", "🏛️")
            pillar_name = pillar_def.get("name", brief.pillar_id or "Other")
            if icon not in pillar_icons:
                pillar_icons.append(icon)
            pillar_counts[pillar_name] = pillar_counts.get(pillar_name, 0) + 1

        # ===== 1. TITLE SLIDE =====
        title_subtitle = (
            f"{' '.join(pillar_icons)} | {len(briefs)} Strategic Trends\n"
            f"{datetime.now(timezone.utc).strftime('%B %Y')}"
        )
        add_title_slide(prs, workstream_name, title_subtitle)

        # ===== 2. PORTFOLIO DASHBOARD =====
        comparison_chart_path = generate_portfolio_comparison_chart(briefs)
        if comparison_chart_path:
            temp_files_to_cleanup.append(comparison_chart_path)

        pillar_chart_path = None
        if pillar_counts:
            pillar_chart_path = generate_pillar_distribution_chart(
                pillar_counts, "Distribution by Pillar"
            )
            if pillar_chart_path:
                temp_files_to_cleanup.append(pillar_chart_path)

        add_portfolio_dashboard_slide(
            prs, briefs, comparison_chart_path, pillar_chart_path
        )

        # ===== 3. WHY THIS MATTERS NOW =====
        urgency = (
            getattr(synthesis, "urgency_statement", "")
            or f"These {len(briefs)} trends represent critical opportunities and challenges. Early action positions Austin as a leader; delay risks falling behind peer cities."
        )
        urgency_content = (
            f"{urgency}\n\n**The Window of Opportunity**\n\n"
            "Cities that move first on emerging trends gain competitive advantage "
            "in talent attraction, federal funding, and citizen satisfaction."
        )
        add_smart_content_slide(
            prs,
            title="Why This Matters Now",
            content=urgency_content,
            max_chars=1200,
        )

        # ===== 4. EXECUTIVE OVERVIEW =====
        overview_content = (
            synthesis.executive_overview or "Portfolio synthesis in progress..."
        )
        add_smart_content_slide(
            prs,
            title="Executive Overview",
            content=overview_content,
            max_chars=1400,
        )

        # ===== 5. PRIORITY MATRIX =====
        if matrix_chart_path := generate_priority_matrix_chart(briefs, synthesis):
            temp_files_to_cleanup.append(matrix_chart_path)

            slide_layout = prs.slide_layouts[6]
            matrix_slide = prs.slides.add_slide(slide_layout)
            add_pptx_header(matrix_slide)
            add_pptx_footer(matrix_slide)

            title_box = matrix_slide.shapes.add_textbox(
                PPTX_MARGIN,
                Inches(1.25),
                PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN),
                Inches(0.5),
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = "Strategic Priority Matrix"
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["primary"])

            try:
                matrix_slide.shapes.add_picture(
                    matrix_chart_path,
                    Inches(0.8),
                    Inches(1.9),
                    width=Inches(8.4),
                    height=Inches(5.0),
                )
            except Exception as e:
                logger.warning(f"Failed to add priority matrix chart: {e}")
        else:
            # Fallback to text-based priority slide
            matrix = synthesis.priority_matrix or {}
            urgent = matrix.get("high_impact_urgent", [])
            strategic = matrix.get("high_impact_strategic", [])
            monitor = matrix.get("monitor", [])

            priority_content = "🔴 **High Impact - Urgent Action**\n" + (
                "\n".join(f"• {item}" for item in urgent)
                if urgent
                else "• None identified"
            )
            priority_content += "\n\n🟡 **High Impact - Strategic Planning**\n"
            priority_content += (
                "\n".join(f"• {item}" for item in strategic)
                if strategic
                else "• None identified"
            )
            priority_content += "\n\n🟢 **Monitor & Evaluate**\n"
            priority_content += (
                "\n".join(f"• {item}" for item in monitor)
                if monitor
                else "• None identified"
            )

            add_smart_content_slide(
                prs,
                title="Strategic Priorities",
                content=priority_content,
                max_chars=1500,
            )

        # ===== 6. IMPLEMENTATION GUIDANCE =====
        impl = getattr(synthesis, "implementation_guidance", {}) or {}
        impl_lines = []
        if impl.get("pilot_now"):
            impl_lines.append(
                f"🚀 **Ready to Pilot**: {', '.join(impl['pilot_now'])}"
            )
        if impl.get("investigate_further"):
            impl_lines.append(
                f"🔍 **Investigate Further**: {', '.join(impl['investigate_further'])}"
            )
        if impl.get("meet_with_vendors"):
            impl_lines.append(
                f"🤝 **Meet with Vendors**: {', '.join(impl['meet_with_vendors'])}"
            )
        if impl.get("policy_review"):
            impl_lines.append(
                f"📋 **Policy Review Needed**: {', '.join(impl['policy_review'])}"
            )
        if impl.get("staff_training"):
            impl_lines.append(
                f"👥 **Staff Training Focus**: {', '.join(impl['staff_training'])}"
            )
        if impl.get("budget_planning"):
            impl_lines.append(
                f"💰 **Budget Planning**: {', '.join(impl['budget_planning'])}"
            )

        if impl_lines:
            impl_content = "What should Austin DO with each trend?\n\n" + "\n".join(
                impl_lines
            )
            add_smart_content_slide(
                prs,
                title="Implementation Guidance",
                content=impl_content,
                max_chars=1400,
            )

        # ===== 7. PER-CARD DEEP DIVES =====
        for i, brief in enumerate(briefs, 1):
            card_chart_path = None
            if (
                (brief.impact_score and brief.impact_score > 0)
                or (brief.relevance_score and brief.relevance_score > 0)
                or (brief.velocity_score and brief.velocity_score > 0)
            ):
                scores = {
                    "Impact": brief.impact_score or 0,
                    "Relevance": brief.relevance_score or 0,
                    "Velocity": brief.velocity_score or 0,
                }
                if valid_scores := {k: v for k, v in scores.items() if v > 0}:
                    card_chart_path = generate_radar_chart(
                        valid_scores, brief.card_name, CHART_DPI
                    )
                    if card_chart_path:
                        temp_files_to_cleanup.append(card_chart_path)

            add_card_deep_dive_slides(prs, brief, i, card_chart_path)

        # ===== 8. CROSS-CUTTING THEMES =====
        themes_content = "**Common Patterns Across Trends**\n"
        themes_content += (
            "\n".join(f"• {theme}" for theme in (synthesis.key_themes or []))
            or "• Analysis in progress"
        )
        themes_content += "\n\n**Strategic Connections**\n"
        themes_content += (
            "\n".join(
                f"• {insight}"
                for insight in (synthesis.cross_cutting_insights or [])
            )
            or "• Analysis in progress"
        )

        add_smart_content_slide(
            prs,
            title="Cross-Cutting Themes",
            content=themes_content,
            max_chars=1400,
        )

        # ===== 9. ACTION PLAN =====
        if ninety_day := getattr(synthesis, "ninety_day_actions", []) or []:
            actions_content = "What Austin should do in the next 90 days:\n\n"
            for action in ninety_day[:5]:
                action_text = action.get("action", "")
                owner = action.get("owner", "TBD")
                by_when = action.get("by_when", "90 days")
                metric = action.get("success_metric", "")
                actions_content += f"✓ **{action_text}**\n"
                actions_content += f"   Owner: {owner} | By: {by_when}"
                if metric:
                    actions_content += f"\n   Success: {metric}"
                actions_content += "\n\n"
        else:
            actions_content = ""
            for action in (synthesis.recommended_actions or [])[:6]:
                action_text = action.get("action", "")
                owner = action.get("owner", "TBD")
                timeline = action.get("timeline", "TBD")
                related_cards = action.get("cards", [])

                actions_content += f"✓ **{action_text}**\n"
                actions_content += f"   Owner: {owner} | Timeline: {timeline}"
                if related_cards:
                    actions_content += f" | Related: {', '.join(related_cards[:2])}"
                actions_content += "\n\n"

            if not actions_content:
                actions_content = (
                    "Action plan to be developed based on leadership priorities."
                )

        add_smart_content_slide(
            prs, title="90-Day Action Plan", content=actions_content, max_chars=1400
        )

        # ===== 10. RISKS & OPPORTUNITIES =====
        risk_text = (
            getattr(synthesis, "risk_summary", "")
            or "Delayed action on these trends could result in Austin falling behind peer cities, missing federal funding windows, and losing competitive advantage."
        )
        opp_text = (
            getattr(synthesis, "opportunity_summary", "")
            or "Early action positions Austin as a national leader, attracts innovation investment, and delivers improved services to residents."
        )

        risk_opp_content = (
            f"⚠️ **If Austin Doesn't Act**\n{risk_text}\n\n"
            f"✨ **If Austin Leads**\n{opp_text}"
        )
        add_smart_content_slide(
            prs,
            title="Risks & Opportunities",
            content=risk_opp_content,
            max_chars=1400,
        )

        # ===== 11. AI DISCLOSURE =====
        add_ai_disclosure_slide(prs)

        temp_file = tempfile.NamedTemporaryFile(
            suffix=".pptx", delete=False, prefix="foresight_portfolio_local_"
        )
        prs.save(temp_file.name)
        temp_file.close()

        logger.info(
            f"Generated enhanced local portfolio PPTX: {len(briefs)} cards, {len(prs.slides)} slides"
        )
        return temp_file.name

    finally:
        cleanup_temp_files(temp_files_to_cleanup)
