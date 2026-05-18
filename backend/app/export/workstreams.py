"""Workstream-level export generators (PDF + PPTX).

Module-level versions of ``ExportService.generate_workstream_pdf`` and
``ExportService.generate_workstream_pptx``. They render a workstream + its
cards into a multi-page PDF or multi-slide PPTX file and return the temp path.

The class methods on ``ExportService`` are now thin facades over these.
"""

import logging
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors as rl_colors
from reportlab.lib.units import inch
from reportlab.platypus import (
    HRFlowable,
    Image as RLImage,
    PageBreak,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
)
from supabase import Client

from ..models.export import CardExportData
from . import charts as _charts
from . import data_access as _data_access
from . import utils as _utils
from .branding import FORESIGHT_COLORS, PDF_COLORS
from .pdf import ProfessionalPDFBuilder, get_professional_pdf_styles
from .pptx import (
    PPTX_MARGIN,
    PPTX_SLIDE_HEIGHT,
    PPTX_SLIDE_WIDTH,
    PPTX_SUBTITLE_FONT_SIZE,
    add_content_slide,
    add_pptx_footer,
    add_pptx_header,
    add_title_slide,
    hex_to_rgb,
)

logger = logging.getLogger(__name__)


# PDF chart dimensions used in this module
PDF_CHART_WIDTH = 5.5 * inch
PDF_CHART_HEIGHT = 4 * inch
PDF_HORIZON_CHART_WIDTH = 4.5 * inch
PDF_HORIZON_CHART_HEIGHT = 3 * inch
PDF_BODY_FONT_SIZE = 11
PDF_SMALL_FONT_SIZE = 9


def _compute_pillar_horizon_counts(
    cards: List[CardExportData],
    horizon_unknown_label: str | None = None,
) -> tuple[Dict[str, int], Dict[str, int]]:
    """Aggregate pillar and horizon counts from a list of cards.

    Pillar uses ``card.pillar_name`` then ``card.pillar_id`` then ``"Unknown"``.
    Horizon defaults to ``card.horizon``; when missing, falls back to
    ``horizon_unknown_label`` if provided, otherwise the card is skipped from
    the horizon bucket. This preserves the PDF (skip empty) vs PPTX (bucket as
    "Unknown") behaviors with a single shared helper.
    """
    pillar_counts: Dict[str, int] = {}
    horizon_counts: Dict[str, int] = {}
    for card in cards:
        pillar = card.pillar_name or card.pillar_id or "Unknown"
        pillar_counts[pillar] = pillar_counts.get(pillar, 0) + 1

        horizon = card.horizon or horizon_unknown_label
        if horizon:
            horizon_counts[horizon] = horizon_counts.get(horizon, 0) + 1

    return pillar_counts, horizon_counts


async def generate_workstream_pdf(
    supabase: Client,
    workstream_id: str,
    include_charts: bool = True,
    max_cards: int = 50,
) -> str:
    """Generate a multi-page workstream PDF and return the temp file path.

    Fetches workstream + cards via ``data_access``. Raises ``ValueError`` if the
    workstream doesn't exist; re-raises any other ReportLab error after logging.
    """
    temp_files = []

    try:
        workstream, cards = await _data_access.get_workstream_cards(
            supabase, workstream_id, max_cards
        )

        if not workstream:
            raise ValueError(f"Workstream {workstream_id} not found")

        pdf_file = tempfile.NamedTemporaryFile(
            suffix=".pdf", delete=False, prefix="foresight_workstream_"
        )
        pdf_path = pdf_file.name
        pdf_file.close()

        styles = get_professional_pdf_styles()

        elements = []

        # Title page
        workstream_name = workstream.get("name", "Workstream Report")
        elements.extend(
            (Paragraph(workstream_name, styles["DocTitle"]), Spacer(1, 12))
        )
        elements.append(
            HRFlowable(
                width="100%",
                thickness=2,
                color=PDF_COLORS["secondary"],
                spaceBefore=6,
                spaceAfter=12,
            )
        )

        # Workstream description
        if workstream.get("description"):
            elements.append(Paragraph("Overview", styles["SectionHeading"]))
            elements.append(
                Paragraph(workstream["description"], styles["BodyText"])
            )
            elements.append(Spacer(1, 12))

        # Summary statistics
        elements.append(Paragraph("Summary", styles["SectionHeading"]))

        if not cards:
            elements.append(
                Paragraph(
                    "No cards currently match this workstream criteria.",
                    styles["BodyText"],
                )
            )
        else:
            summary_text = (
                f"This workstream contains <b>{len(cards)}</b> intelligence cards."
            )
            if len(cards) >= max_cards:
                summary_text += f" (Showing first {max_cards})"
            elements.append(Paragraph(summary_text, styles["BodyText"]))
            elements.append(Spacer(1, 12))

            # Distribution charts
            if include_charts and cards:
                pillar_counts, horizon_counts = _compute_pillar_horizon_counts(
                    cards
                )

                if pillar_chart_path := _charts.generate_pillar_distribution_chart(
                    pillar_counts
                ):
                    temp_files.append(pillar_chart_path)
                    try:
                        img = RLImage(
                            pillar_chart_path,
                            width=PDF_CHART_WIDTH,
                            height=PDF_CHART_HEIGHT,
                        )
                        elements.append(img)
                        elements.append(Spacer(1, 12))
                    except Exception as e:
                        logger.warning(f"Failed to add pillar chart to PDF: {e}")

                if horizon_chart_path := _charts.generate_horizon_distribution_chart(
                    horizon_counts
                ):
                    temp_files.append(horizon_chart_path)
                    try:
                        img = RLImage(
                            horizon_chart_path,
                            width=PDF_HORIZON_CHART_WIDTH,
                            height=PDF_HORIZON_CHART_HEIGHT,
                        )
                        elements.append(img)
                        elements.append(Spacer(1, 12))
                    except Exception as e:
                        logger.warning(f"Failed to add horizon chart to PDF: {e}")

            # Cards table summary
            elements.append(PageBreak())
            elements.append(Paragraph("Cards Overview", styles["SectionHeading"]))

            table_data = [["Name", "Pillar", "Horizon", "Impact"]]
            for card in cards:
                table_data.append(
                    [
                        card.name[:40] + ("..." if len(card.name) > 40 else ""),
                        card.pillar_name or card.pillar_id or "N/A",
                        card.horizon or "N/A",
                        _utils.format_score_display(card.impact_score),
                    ]
                )

            table = Table(
                table_data,
                colWidths=[2.5 * inch, 1.5 * inch, 0.8 * inch, 0.8 * inch],
            )
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), PDF_COLORS["primary"]),
                        ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("FONTSIZE", (0, 0), (-1, 0), PDF_BODY_FONT_SIZE),
                        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                        ("FONTSIZE", (0, 1), (-1, -1), PDF_SMALL_FONT_SIZE),
                        ("ALIGN", (2, 0), (-1, -1), "CENTER"),
                        (
                            "ROWBACKGROUNDS",
                            (0, 1),
                            (-1, -1),
                            [rl_colors.white, PDF_COLORS["light"]],
                        ),
                        ("GRID", (0, 0), (-1, -1), 0.5, PDF_COLORS["light"]),
                        ("BOX", (0, 0), (-1, -1), 1, PDF_COLORS["primary"]),
                        ("TOPPADDING", (0, 0), (-1, -1), 6),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                        ("LEFTPADDING", (0, 0), (-1, -1), 6),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ]
                )
            )
            elements.append(table)

            # Individual card details (paginated)
            elements.append(PageBreak())
            elements.append(Paragraph("Card Details", styles["SectionHeading"]))
            elements.append(Spacer(1, 12))

            for i, card in enumerate(cards):
                if i > 0:
                    elements.append(Spacer(1, 18))
                    elements.append(
                        HRFlowable(
                            width="80%",
                            thickness=0.5,
                            color=PDF_COLORS["light"],
                            spaceBefore=6,
                            spaceAfter=12,
                        )
                    )

                elements.append(Paragraph(card.name, styles["SubsectionHeading"]))

                badge_parts = []
                if card.pillar_name or card.pillar_id:
                    badge_parts.append(
                        f"<b>Pillar:</b> {card.pillar_name or card.pillar_id}"
                    )
                if card.horizon:
                    badge_parts.append(f"<b>Horizon:</b> {card.horizon}")
                if card.stage_name or card.stage_id:
                    badge_parts.append(
                        f"<b>Stage:</b> {card.stage_name or card.stage_id}"
                    )
                if badge_parts:
                    elements.append(
                        Paragraph(" | ".join(badge_parts), styles["SmallText"])
                    )
                    elements.append(Spacer(1, 6))

                if card.summary:
                    elements.append(Paragraph(card.summary, styles["BodyText"]))

                key_scores = []
                if card.impact_score is not None:
                    key_scores.append(f"Impact: {card.impact_score}")
                if card.relevance_score is not None:
                    key_scores.append(f"Relevance: {card.relevance_score}")
                if card.maturity_score is not None:
                    key_scores.append(f"Maturity: {card.maturity_score}")
                if key_scores:
                    elements.append(
                        Paragraph(
                            "<i>Scores: " + " | ".join(key_scores) + "</i>",
                            styles["SmallText"],
                        )
                    )

        elements.append(Spacer(1, 24))
        elements.append(
            HRFlowable(
                width="100%",
                thickness=1,
                color=PDF_COLORS["light"],
                spaceBefore=6,
                spaceAfter=6,
            )
        )

        footer_text = (
            f"Export Date: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}"
        )
        elements.append(Paragraph(footer_text, styles["SmallText"]))

        builder = ProfessionalPDFBuilder(
            filename=pdf_path,
            title=workstream_name,
            include_logo=True,
            include_ai_disclosure=True,
        )
        builder.build(elements)

        logger.info(
            f"Generated professional workstream PDF for: {workstream_name} "
            f"with {len(cards)} cards"
        )
        return pdf_path

    except Exception as e:
        logger.error(f"Error generating workstream PDF {workstream_id}: {e}")
        raise

    finally:
        _utils.cleanup_temp_files(temp_files)


async def generate_workstream_pptx(
    workstream: Dict[str, Any],
    cards: List[CardExportData],
    include_charts: bool = True,
    include_card_details: bool = True,
) -> str:
    """Generate a workstream PPTX (title / summary / distribution / per-card slides).

    Returns the temp file path.
    """
    temp_files_to_cleanup = []

    try:
        workstream_name = workstream.get("name", "Workstream Report")
        logger.info(f"Generating workstream PowerPoint: {workstream_name}")

        prs = Presentation()
        prs.slide_width = PPTX_SLIDE_WIDTH
        prs.slide_height = PPTX_SLIDE_HEIGHT

        # 1. Title slide
        add_title_slide(
            prs,
            title=workstream_name,
            subtitle=f"Intelligence Report • {len(cards)} Cards",
        )

        # 2. Summary slide
        summary_items = [
            ("Total Cards", str(len(cards))),
            ("Description", workstream.get("description", "N/A")),
        ]

        pillar_counts, horizon_counts = _compute_pillar_horizon_counts(
            cards, horizon_unknown_label="Unknown"
        )

        if pillar_counts:
            pillar_summary = ", ".join(
                f"{k}: {v}" for k, v in pillar_counts.items()
            )
            summary_items.append(("Pillars", pillar_summary))

        if horizon_counts:
            horizon_summary = ", ".join(
                f"{k}: {v}" for k, v in horizon_counts.items()
            )
            summary_items.append(("Horizons", horizon_summary))

        add_content_slide(
            prs, title="Workstream Summary", content_items=summary_items
        )

        # 3. Distribution charts slide
        if include_charts and cards:
            pillar_chart_path = None
            if pillar_counts:
                pillar_chart_path = _charts.generate_pillar_distribution_chart(
                    pillar_counts
                )
                if pillar_chart_path:
                    temp_files_to_cleanup.append(pillar_chart_path)

            horizon_chart_path = None
            if horizon_counts:
                horizon_chart_path = _charts.generate_horizon_distribution_chart(
                    horizon_counts
                )
                if horizon_chart_path:
                    temp_files_to_cleanup.append(horizon_chart_path)

            if pillar_chart_path or horizon_chart_path:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)

                add_pptx_header(slide)
                add_pptx_footer(slide)

                title_box = slide.shapes.add_textbox(
                    PPTX_MARGIN,
                    Inches(1.25),
                    PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN),
                    Inches(0.6),
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = "Distribution Analysis"
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                title_para.font.color.rgb = hex_to_rgb(
                    FORESIGHT_COLORS["primary"]
                )

            if pillar_chart_path and Path(pillar_chart_path).exists():
                try:
                    slide.shapes.add_picture(
                        pillar_chart_path,
                        Inches(0.3),
                        Inches(2.0),
                        width=Inches(5.5),
                        height=Inches(4.0),
                    )
                except Exception as e:
                    logger.warning(f"Failed to add pillar chart: {e}")

            if horizon_chart_path and Path(horizon_chart_path).exists():
                try:
                    slide.shapes.add_picture(
                        horizon_chart_path,
                        Inches(6.5),
                        Inches(2.0),
                        width=Inches(5.5),
                        height=Inches(4.0),
                    )
                except Exception as e:
                    logger.warning(f"Failed to add horizon chart: {e}")

        if include_card_details:
            if cards:
                for card in cards[:50]:
                    card_items = [
                        ("Summary", card.summary),
                        ("Pillar", card.pillar_name or card.pillar_id),
                        ("Horizon", card.horizon),
                        ("Stage", card.stage_name or card.stage_id),
                    ]
                    scores = card.get_all_scores()
                    if valid_scores := {
                        k: v for k, v in scores.items() if v is not None
                    }:
                        scores_text = ", ".join(
                            f"{k}: {v}" for k, v in valid_scores.items()
                        )
                        card_items.append(("Scores", scores_text))

                    card_items = [(k, v) for k, v in card_items if v]

                    add_content_slide(
                        prs, title=card.name[:50], content_items=card_items
                    )

            else:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)

                add_pptx_header(slide)
                add_pptx_footer(slide)

                msg_box = slide.shapes.add_textbox(
                    Inches(2), Inches(3.5), Inches(9), Inches(2)
                )
                msg_frame = msg_box.text_frame
                msg_para = msg_frame.paragraphs[0]
                msg_para.text = "No cards currently match this workstream criteria"
                msg_para.font.size = PPTX_SUBTITLE_FONT_SIZE
                msg_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["dark"])
                msg_para.alignment = PP_ALIGN.CENTER

        temp_file = tempfile.NamedTemporaryFile(
            suffix=".pptx", delete=False, prefix="foresight_workstream_"
        )
        prs.save(temp_file.name)

        logger.info(f"Workstream PowerPoint generated: {temp_file.name}")
        return temp_file.name

    except Exception as e:
        logger.error(f"Error generating workstream PowerPoint: {e}")
        raise

    finally:
        _utils.cleanup_temp_files(temp_files_to_cleanup)
