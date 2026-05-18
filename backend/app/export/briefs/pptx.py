"""PowerPoint generators for executive briefs.

Tries Gamma.app for AI-powered presentation generation first; falls back to
local python-pptx rendering with classification appendix slides.
"""

import logging
import re
import tempfile
from datetime import datetime
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt

from app.openai_provider import get_chat_deployment

from ..branding import FORESIGHT_COLORS
from ..pptx import (
    PPTX_MARGIN,
    PPTX_SLIDE_HEIGHT,
    PPTX_SLIDE_WIDTH,
    add_pptx_footer,
    add_pptx_header,
    add_title_slide,
    hex_to_rgb,
)

logger = logging.getLogger(__name__)


def clean_markdown_for_pptx(text: str) -> str:
    """Clean markdown text for PowerPoint display.

    Removes markdown artifacts that don't render well in PPTX.
    """
    if not text:
        return ""

    text = re.sub(r"```[\s\S]*?```", "", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)

    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"__(.+?)__", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"_(.+?)_", r"\1", text)

    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)

    text = re.sub(r"^[\-\*•]\s+", "• ", text, flags=re.MULTILINE)
    text = re.sub(r"^[-*_]{3,}\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)

    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"  +", " ", text)

    return text.strip()


def parse_markdown_sections_improved(
    content_markdown: str, max_sections: int = 10
) -> List[Tuple[str, str]]:
    """Improved markdown section parser.

    Handles:
    - Multiple header formats (#, ##, ###)
    - ALL CAPS section headers
    - Bold section headers (**Section**)
    - Reasonable content length per section
    """
    sections: List[Tuple[str, str]] = []
    current_title = "Overview"
    current_content: List[str] = []

    lines = content_markdown.split("\n")

    for line in lines:
        line_stripped = line.strip()

        if header_match := re.match(r"^(#{1,3})\s+(.+)$", line_stripped):
            if current_content:
                content_text = "\n".join(current_content).strip()
                if content_text and len(content_text) > 30:
                    sections.append((current_title, content_text))
            current_title = header_match.group(2).strip()
            current_title = re.sub(r"\s*#+\s*$", "", current_title)
            current_title = re.sub(r"^\*\*|\*\*$", "", current_title)
            current_content = []
            continue

        if re.match(
            r"^[A-Z][A-Z\s&\-]{5,}$", line_stripped
        ) and not line_stripped.startswith("•"):
            if current_content:
                content_text = "\n".join(current_content).strip()
                if content_text and len(content_text) > 30:
                    sections.append((current_title, content_text))
            current_title = line_stripped.title()
            current_content = []
            continue

        if bold_match := re.match(r"^\*\*([^*]+)\*\*:?\s*$", line_stripped):
            if current_content:
                content_text = "\n".join(current_content).strip()
                if content_text and len(content_text) > 30:
                    sections.append((current_title, content_text))
            current_title = bold_match.group(1).strip()
            current_content = []
            continue

        current_content.append(line)

    if current_content:
        content_text = "\n".join(current_content).strip()
        if content_text and len(content_text) > 30:
            sections.append((current_title, content_text))

    if not sections and content_markdown.strip():
        sections = [("Key Findings", content_markdown.strip())]

    return sections[:max_sections]


def add_smart_content_slide(
    prs: Presentation, title: str, content: str, max_chars: int = 1000
) -> None:
    """Add a content slide with smart text handling.

    Handles bullet points, truncation, and proper formatting.
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    add_pptx_header(slide)
    add_pptx_footer(slide)

    title_box = slide.shapes.add_textbox(
        PPTX_MARGIN, Inches(1.25), PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title[:60] if len(title) > 60 else title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["primary"])

    if len(content) > max_chars:
        content = f"{content[:max_chars - 3]}..."

    content_box = slide.shapes.add_textbox(
        PPTX_MARGIN, Inches(1.95), PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(4.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    lines = content.split("\n")
    first_para = True

    for line in lines:
        line = line.strip()
        if not line:
            continue

        if first_para:
            para = content_frame.paragraphs[0]
            first_para = False
        else:
            para = content_frame.add_paragraph()

        if line.startswith("•") or line.startswith("-") or line.startswith("*"):
            line = line.lstrip("•-* ")
            para.text = f"• {line}"
            para.level = 0
        else:
            para.text = line

        para.font.size = Pt(16)
        para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["dark"])
        para.space_before = Pt(6)
        para.space_after = Pt(4)


def add_ai_disclosure_slide(prs: Presentation) -> None:
    """Add an AI technology disclosure slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_pptx_header(slide)
    add_pptx_footer(slide)

    title_box = slide.shapes.add_textbox(
        PPTX_MARGIN, Inches(1.25), PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "About This Report"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["primary"])

    disclosure_text = f"""This strategic intelligence brief was generated using the FORESIGHT platform, powered by advanced AI technologies:

• OpenAI {get_chat_deployment()} - Strategic analysis, classification, and scoring
• GPT Researcher - Autonomous deep research orchestration
• SearXNG and Serper - Web search aggregation
• trafilatura - Article extraction from source URLs

The City of Austin is committed to transparent and responsible use of AI technology in public service. All AI-generated content is reviewed for accuracy and relevance."""

    content_box = slide.shapes.add_textbox(
        PPTX_MARGIN, Inches(1.95), PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(4.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    lines = disclosure_text.split("\n")
    first_para = True

    for line in lines:
        line = line.strip()
        if not line:
            continue

        if first_para:
            para = content_frame.paragraphs[0]
            first_para = False
        else:
            para = content_frame.add_paragraph()

        para.text = line
        para.font.size = Pt(14)
        para.font.color.rgb = hex_to_rgb(FORESIGHT_COLORS["dark"])
        para.space_before = Pt(4)
        para.space_after = Pt(4)


async def generate_brief_pptx(
    brief_title: str,
    card_name: str,
    executive_summary: str,
    content_markdown: str,
    generated_at: Optional[datetime] = None,
    version: Optional[int] = None,
    classification: Optional[Dict[str, str]] = None,
    use_gamma: bool = True,
) -> str:
    """Generate a PowerPoint presentation for an executive brief.

    Attempts to use Gamma.app for AI-powered presentation generation.
    Falls back to local python-pptx generation if Gamma is unavailable.
    """
    if use_gamma:
        try:
            from app.gamma_service import GammaService

            gamma = GammaService()
            if gamma.is_available():
                logger.info(
                    f"Attempting Gamma.app presentation generation for: {brief_title}"
                )

                result = await gamma.generate_presentation(
                    title=brief_title,
                    executive_summary=executive_summary,
                    content_markdown=content_markdown,
                    classification=classification,
                    num_slides=8,
                    include_images=True,
                    export_format="pptx",
                )

                if result.success and result.pptx_url:
                    pptx_bytes = await gamma.download_export(result.pptx_url)

                    if pptx_bytes:
                        temp_file = tempfile.NamedTemporaryFile(
                            suffix=".pptx",
                            delete=False,
                            prefix="foresight_gamma_brief_",
                        )
                        temp_file.write(pptx_bytes)
                        temp_file.close()

                        logger.info(
                            f"Gamma presentation generated successfully: {temp_file.name}"
                        )
                        if result.credits_used:
                            logger.info(
                                f"Gamma credits used: {result.credits_used}, remaining: {result.credits_remaining}"
                            )

                        return temp_file.name

                if result.error_message:
                    logger.warning(
                        f"Gamma generation failed: {result.error_message}"
                    )
                else:
                    logger.warning(
                        "Gamma generation completed but no PPTX URL returned"
                    )

        except ImportError:
            logger.warning("Gamma service not available, using fallback")
        except Exception as e:
            logger.warning(f"Gamma generation error, falling back to local: {e}")

    return await generate_brief_pptx_local(
        brief_title=brief_title,
        card_name=card_name,
        executive_summary=executive_summary,
        content_markdown=content_markdown,
        generated_at=generated_at,
        version=version,
        classification=classification,
    )


async def generate_brief_pptx_local(
    brief_title: str,
    card_name: str,
    executive_summary: str,
    content_markdown: str,
    generated_at: Optional[datetime] = None,
    version: Optional[int] = None,
    classification: Optional[Dict[str, str]] = None,
) -> str:
    """Local PowerPoint generation with improved markdown handling.

    Used as fallback when Gamma.app is unavailable. Produces:
    - Title slide with classification tags
    - Executive summary slide
    - Up to 8 content slides
    - AI disclosure slide
    - Appendix slides defining each classification tag
    """
    from app.gamma_service import (
        HORIZON_DEFINITIONS,
        HORIZON_NAMES,
        PILLAR_DEFINITIONS,
        PILLAR_NAMES,
        STAGE_DEFINITIONS,
        STAGE_NAMES,
    )

    try:
        logger.info(f"Generating local brief PowerPoint: {brief_title}")

        prs = Presentation()
        prs.slide_width = PPTX_SLIDE_WIDTH
        prs.slide_height = PPTX_SLIDE_HEIGHT

        used_pillar = None
        used_horizon = None
        used_stage = None

        subtitle = "Strategic Intelligence Brief"
        if classification:
            tag_parts = []
            if classification.get("pillar"):
                pillar = classification["pillar"].upper()
                pillar_def = PILLAR_DEFINITIONS.get(pillar, {})
                pillar_name = pillar_def.get(
                    "name", PILLAR_NAMES.get(pillar, pillar)
                )
                pillar_icon = pillar_def.get("icon", "")
                tag_parts.append(f"{pillar_icon} {pillar_name}")
                used_pillar = pillar
            if classification.get("horizon"):
                horizon = classification["horizon"].upper()
                horizon_def = HORIZON_DEFINITIONS.get(horizon, {})
                horizon_name = horizon_def.get(
                    "name", HORIZON_NAMES.get(horizon, horizon)
                )
                horizon_icon = horizon_def.get("icon", "")
                tag_parts.append(f"{horizon_icon} {horizon_name}")
                used_horizon = horizon
            if classification.get("stage"):
                stage_raw = classification["stage"]
                if stage_match := re.search(r"(\d+)", str(stage_raw)):
                    stage_num = int(stage_match.group(1))
                    stage_def = STAGE_DEFINITIONS.get(stage_num, {})
                    stage_name = stage_def.get(
                        "name", STAGE_NAMES.get(stage_num, f"Stage {stage_num}")
                    )
                    stage_icon = stage_def.get("icon", "")
                    tag_parts.append(
                        f"{stage_icon} Stage {stage_num}: {stage_name}"
                    )
                    used_stage = stage_num
            if tag_parts:
                subtitle = "  |  ".join(tag_parts)

        if generated_at:
            subtitle += f"\n{generated_at.strftime('%B %d, %Y')}"

        add_title_slide(prs, title=brief_title, subtitle=subtitle)

        if executive_summary:
            clean_summary = clean_markdown_for_pptx(executive_summary)
            add_smart_content_slide(
                prs,
                title="Executive Summary",
                content=clean_summary,
                max_chars=1200,
            )

        if content_markdown:
            sections = parse_markdown_sections_improved(content_markdown)

            for section_title, section_content in sections[:8]:
                clean_content = clean_markdown_for_pptx(section_content)
                add_smart_content_slide(
                    prs,
                    title=section_title,
                    content=clean_content,
                    max_chars=1000,
                )

        add_ai_disclosure_slide(prs)

        appendix_content = """The following slides provide context for the strategic classification tags used in this brief.

These definitions help ensure consistent understanding across City departments and leadership."""
        add_smart_content_slide(
            prs,
            title="Appendix: Classification Reference",
            content=appendix_content,
            max_chars=1000,
        )

        if used_pillar and used_pillar in PILLAR_DEFINITIONS:
            pillar_def = PILLAR_DEFINITIONS[used_pillar]
            pillar_content = f"""Definition: {pillar_def['description']}

Focus Areas:
"""
            for area in pillar_def.get("focus_areas", []):
                pillar_content += f"- {area}\n"
            pillar_content += """
This pillar is one of six strategic focus areas guiding City of Austin planning and investment decisions."""

            add_smart_content_slide(
                prs,
                title=f"{pillar_def.get('icon', '')} Strategic Pillar: {pillar_def['name']}",
                content=pillar_content,
                max_chars=1500,
            )

        if used_horizon and used_horizon in HORIZON_DEFINITIONS:
            horizon_def = HORIZON_DEFINITIONS[used_horizon]
            horizon_content = f"""Timeframe: {horizon_def['timeframe']}

Definition: {horizon_def['description']}

Characteristics:
"""
            for char in horizon_def.get("characteristics", []):
                horizon_content += f"- {char}\n"
            horizon_content += """
The planning horizon indicates when this trend is expected to require significant City attention or action."""

            add_smart_content_slide(
                prs,
                title=f"{horizon_def.get('icon', '')} Planning Horizon: {horizon_def['name']}",
                content=horizon_content,
                max_chars=1500,
            )

        if used_stage and used_stage in STAGE_DEFINITIONS:
            stage_def = STAGE_DEFINITIONS[used_stage]
            stage_content = f"""Definition: {stage_def['description']}

Key Indicators:
"""
            for indicator in stage_def.get("indicators", []):
                stage_content += f"- {indicator}\n"
            stage_content += """
The maturity stage reflects the current development status of this trend and helps inform appropriate City response strategies."""

            add_smart_content_slide(
                prs,
                title=f"{stage_def.get('icon', '')} Maturity Stage {used_stage}: {stage_def['name']}",
                content=stage_content,
                max_chars=1500,
            )

        temp_file = tempfile.NamedTemporaryFile(
            suffix=".pptx", delete=False, prefix="foresight_brief_"
        )
        prs.save(temp_file.name)

        logger.info(f"Local brief PowerPoint generated: {temp_file.name}")
        return temp_file.name

    except Exception as e:
        logger.error(f"Error generating local brief PowerPoint: {e}")
        raise
